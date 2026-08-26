#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Презентация-КОПИЯ утверждённой руководством колоды (БРУНО_исправленно.pptx),
собранная из model.json.gz конвейера 2.0.

    python build_report_v5.py --model model.json.gz --out Bruno_отчёт_ВТБ.pptx

Вёрстка, палитра, порядок и состав слайдов перенесены из pptx_pipeline/*
(style.py + charts.py + table_data.py + build_pptx.py) без изменений.
Согласованные отличия от оригинала — ровно три:

  1. На слайде уборки к графику плановых задач добавлен парный график
     НЕПЛАНОВЫХ задач по уборке; табличка скорости по башням переехала вниз.
  2. Списки объектов и комендантов в приложениях не фиксированы (жёсткие
     APPENDIX_OBJECT_IDS / EXCLUDED_KOMENDANT_IDS из style.py больше не
     применяются), а состоят из активных в ТЕКУЩЕМ месяце.
  3. Горизонт — текущий (неполный) месяц + 3 предыдущих.

Слайд «Внедрение: оклейка QR» остался с июльским содержимым: этих данных нет
в API Bruno, они ведутся вручную. По решению заказчика слайд оставлен как есть
для ручной правки в готовом файле (см. IMPLEMENTATION ниже — правится в одном
месте).

ПРОХОД 17 (замечания заказчика 26.08.2026):
  1. Кегль заголовков и значений KPI подбирается под ширину (fit_size) — на
     слайдах 6/7/8 длинный заголовок и пара «всего/выполнено» переносились на
     вторую строку и наезжали на подпись.
  2. Слайд «Уборка»: вместо одной колонки «Задач/чел.» четыре — «Назначено на
     чел.», «Выполнено на чел.», «Мин./зад. среднее», «Мин./зад. медиана».
  3. Слайд «Итоги периода»: эффективность уборки стала парой
     «плановая / неплановая».
  4. Мини-графики приложений: 3x3 на слайд (было до 4x4) с переносом на
     следующий слайд, обе шкалы подписаны в каждой ячейке, шкала эффективности
     подстраивается под объект, месяцы разделены линией и «рисочкой», под
     каждым месяцем — его эффективность.
  5. Строки приложений отбираются по ПОБЛОЧНОЙ активности в текущем месяце, а
     задачи «в работе» (NEW/WAITING/COMPLETING) показаны отдельно — жёлтым
     сегментом и третьим числом в таблицах.
  6. Таблица «Все задачи по объектам»: каждый месяц занимает две колонки —
     объём и эффективность месяца; добавлены такие же таблицы отдельно по
     плановым и по неплановым задачам.

ВНИМАНИЕ: пункт 2 требует модель, собранную build_model.py прохода 17 и новее
(в кубе cleanSpeedDay появились поля "all" и "factHist"). На старой модели
колонки «Назначено на чел.» и «медиана» будут прочерками.
"""

import argparse
import datetime as dt
import math
import os
import sys
import tempfile
import textwrap

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import report_v5_model as M

# ============================================================== палитра / стиль
NAVY = "#002882"
NAVY_DARK = "#001A57"
BLUE = "#00AAFF"
GREEN = "#23A75B"
RED = "#E62632"
AMBER = "#F5A623"
WHITE = "#FFFFFF"
BG_LIGHT = "#F4F6FA"
GRID = "#E2E7F0"
TEXT_DARK = "#1A2233"
TEXT_MUTED = "#5B6B82"
TEXT_ON_NAVY_MUTED = "#A9BEDD"

CATEGORICAL_16 = [
    "#002882", "#00AAFF", "#23A75B", "#F5A623", "#E62632", "#6A4C9C",
    "#0E9488", "#C9922C", "#5B6B82", "#B0559C", "#3F8F5F", "#7B8FCB",
    "#A5541F", "#4FA8C9", "#8C1F3B", "#9AA5B1",
]

FONT = "Golos Text"            # корпоративный шрифт ВТБ (для текста в PPTX)
FONT_MPL = "DejaVu Sans"       # для картинок графиков

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5
MARGIN_IN = 0.5
DPI = 200

MONTHS_SHORT = ["", "янв", "фев", "мар", "апр", "май", "июн", "июл",
                "авг", "сен", "окт", "ноя", "дек"]

# --- содержимое слайда «Внедрение: оклейка QR» ---------------------------
# Данных нет в API Bruno, они ведутся вручную. Списки и сноска перенесены из
# утверждённой колоды дословно (включая её орфографию) — правится здесь.
IMPLEMENTATION = [
    ("Готово", ["БФЗ", "БФВ", "БЕ", "Магистральная, 25А", "Трубная, 2",
                "3-я Рыбинская, 18 стр 5", "2-ой Южнопортовый, 18 стр 1"], GREEN, None),
    ("На оклейке", ["Воронцовская, 43А", "Воронцовская, 43Б"], AMBER, "до 31.07.2026"),
    ("В типографии", ["Воздвиженка, 9", "Зотов «Центр» (50%)", "Пакгаузное, 1",
                      "Перовское, 1", "Днепропетровский, 4А"], TEXT_MUTED, "до 24.07.2026"),
]
IMPLEMENTATION_NOTE = "*реализованны за последний месяц"

# Порог для персональных строк/мини-графиков по комендантам в приложениях
# (решение заказчика 22.08.2026). Коменданты с единичными задачами за период
# давали пустые графики и строки вида «1/0 · 100%». В ОБЩИЕ цифры на слайде
# они входят по-прежнему — отсекается только персональная детализация.
KOM_MIN_TASKS = 10
IMPLEMENTATION_TOTAL = 14   # всего зданий в плане оклейки (для слайда «Итоги»)


def band_ylim(values, frac_lo, frac_hi):
    """Границы оси так, чтобы РЯД values занял по высоте ровно полосу
    [frac_lo, frac_hi] — независимо от того, какие в нём числа. Нужно, чтобы
    линия эффективности всегда шла над столбцами и не задевала их подписи."""
    vals = [v for v in values if v is not None]
    if not vals:
        return 0.0, 1.0
    vmin, vmax = min(vals), max(vals)
    if vmax == vmin:                      # ряд-константа: даём ему искусственный размах
        pad = abs(vmax) * 0.1 or 1.0
        vmin, vmax = vmax - pad, vmax + pad
    height = (vmax - vmin) / (frac_hi - frac_lo)
    lo = vmin - frac_lo * height
    return lo, lo + height


def fit_size(text, width_in, base, min_size=9.0, bold=False, factor=None):
    """Кегль, при котором самая длинная строка ГАРАНТИРОВАННО влезает в ширину.

    Зачем (проход 17). python-pptx не умеет autofit, а текстбокс с word_wrap
    просто растёт вниз — в PDF из LibreOffice это видно как наезд второй строки
    заголовка на подзаголовок (слайды 6, 7, 8) и значения KPI на его подпись
    («47 122/ 8 308»). Ширину глифа оцениваем сверху: у Golos Text цифры и
    прописные занимают ~0,55 em, поэтому запас есть."""
    lines = [l for l in str(text).split("\n") if l]
    if not lines:
        return base
    f = factor if factor is not None else (0.56 if bold else 0.52)
    longest = max(len(l) for l in lines)
    fits = width_in * 72.0 / (f * longest)
    return round(max(min_size, min(base, fits)), 1)


def cat_color(i):
    return CATEGORICAL_16[i % len(CATEGORICAL_16)]


def ru(s):
    return s.replace(".", ",")


def fmt_n(v):
    return f"{v:,.0f}".replace(",", " ")


# ================================================================== ГРАФИКИ
class Charts:
    """Всё то же, что в pptx_pipeline/charts.py: PNG содержит только сам график
    (столбцы/линии/шкалу времени), а легенда, подписи осей и заголовки рисуются
    нативными текстовыми блоками PPTX поверх картинки."""

    def __init__(self, outdir):
        self.dir = outdir
        os.makedirs(outdir, exist_ok=True)
        plt.rcParams["font.family"] = FONT_MPL
        plt.rcParams["axes.unicode_minus"] = False

    def _save(self, fig, name, tight=True):
        path = os.path.join(self.dir, name)
        if tight:
            fig.savefig(path, dpi=DPI, bbox_inches="tight", facecolor="white")
        else:
            fig.savefig(path, dpi=DPI, facecolor="white")
        plt.close(fig)
        return path

    @staticmethod
    def _clean_ax(ax):
        for sp in ("top", "right"):
            ax.spines[sp].set_visible(False)
        for sp in ("left", "bottom"):
            ax.spines[sp].set_color(GRID)
        ax.tick_params(colors=TEXT_MUTED, labelsize=8)
        ax.set_axisbelow(True)

    @staticmethod
    def _timeline(ax, weeks, fontsize=7.5, month_fontsize=9.5, label_every=1):
        n = len(weeks)
        ax.set_xlim(-0.6, n - 0.4)
        shown = [i for i in range(n) if i % label_every == 0]
        ax.set_xticks(shown)
        ax.set_xticklabels([weeks[i]["label"] for i in shown], fontsize=fontsize, color=TEXT_MUTED)
        groups, cur, start = [], (weeks[0]["year"], weeks[0]["month"]), 0
        for i in range(1, n + 1):
            key = (weeks[i]["year"], weeks[i]["month"]) if i < n else None
            if key != cur:
                groups.append((start, i - 1, cur))
                if i < n:
                    cur, start = key, i
        for (s, e, (y, m)) in groups:
            ax.annotate(MONTHS_SHORT[m].capitalize(), xy=((s + e) / 2, 0),
                        xycoords=("data", "axes fraction"), xytext=(0, -24),
                        textcoords="offset points", ha="center", va="top",
                        fontsize=month_fontsize, fontweight="bold", color=NAVY,
                        annotation_clip=False)
            if s > 0:
                ax.axvline(s - 0.5, color=GRID, linewidth=1.0, zorder=0)
        ax.tick_params(axis="x", length=0)

    # --- 1. stacked area: обращения по объектам ---
    def feedback_area(self, model, name="b2_main.png", figsize=(15.6, 3.9)):
        weeks, b2 = model["weeks"], model["block2"]
        order = b2["objects_sorted"]
        fig, ax = plt.subplots(figsize=figsize)
        if order:
            x = np.arange(len(weeks))
            series = [np.array(b2["obj_week"][oid]) for oid in order]
            ax.stackplot(x, *series, colors=[cat_color(i) for i in range(len(order))],
                         edgecolor="white", linewidth=0.4)
        self._clean_ax(ax)
        ax.grid(axis="y", color=GRID, linewidth=0.8)
        self._timeline(ax, weeks)
        fig.subplots_adjust(bottom=0.16)
        return self._save(fig, name)

    # --- 2. combo: столбцы вып/проп/в работе + линия эффективности ---
    def combo(self, weeks, completed, missed, completing, efficiency, name, figsize=(15.6, 3.9)):
        x = np.arange(len(weeks))
        fig, ax = plt.subplots(figsize=figsize)
        ax.bar(x, completed, color=GREEN, width=0.62, zorder=3)
        ax.bar(x, missed, bottom=completed, color=RED, width=0.62, zorder=3)
        has_completing = any(c > 0 for c in completing)
        if has_completing:
            ax.bar(x, completing, bottom=np.array(completed) + np.array(missed),
                   color=AMBER, width=0.62, zorder=3)
        self._clean_ax(ax)
        ax.grid(axis="y", color=GRID, linewidth=0.8, zorder=0)
        ax2 = ax.twinx()
        xs = [i for i, v in enumerate(efficiency) if v is not None]
        ys = [v for v in efficiency if v is not None]
        ax2.plot(xs, ys, color=NAVY, linewidth=2.2, marker="o", markersize=4, zorder=4)
        ax2.set_ylim(0, 100)
        ax2.tick_params(colors=NAVY, labelsize=8)
        ax2.spines["top"].set_visible(False)
        for sp in ("left", "bottom", "right"):
            ax2.spines[sp].set_visible(sp == "right")
        ax2.spines["right"].set_color(NAVY)
        self._timeline(ax, weeks)
        fig.subplots_adjust(bottom=0.18)
        return self._save(fig, name), has_completing

    # --- 3. two-part: выполнено / не выполнено + линия эффективности ---
    def two_part(self, weeks, completed, not_completed, efficiency, name, figsize=(15.6, 3.9)):
        x = np.arange(len(weeks))
        fig, ax = plt.subplots(figsize=figsize)
        ax.bar(x, completed, color=GREEN, width=0.62, zorder=3)
        ax.bar(x, not_completed, bottom=completed, color=RED, width=0.62, zorder=3)
        self._clean_ax(ax)
        ax.grid(axis="y", color=GRID, linewidth=0.8, zorder=0)
        ax2 = ax.twinx()
        xs = [i for i, v in enumerate(efficiency) if v is not None]
        ys = [v for v in efficiency if v is not None]
        ax2.plot(xs, ys, color=NAVY, linewidth=2.0, marker="o", markersize=3.5, zorder=4)
        ax2.set_ylim(0, 100)
        ax2.tick_params(colors=NAVY, labelsize=8)
        for sp in ("left", "bottom", "right"):
            ax2.spines[sp].set_visible(sp == "right")
        ax2.spines["right"].set_color(NAVY)
        ax2.spines["top"].set_visible(False)
        self._timeline(ax, weeks)
        fig.subplots_adjust(bottom=0.18)
        return self._save(fig, name)

    # --- 4. сетка мини-графиков (приложения) ---
    # Оформление согласовано 26.08.2026 (вариант «Б»):
    #   * шкала количества (слева) и шкала эффективности (справа) подписаны в
    #     КАЖДОЙ ячейке — раньше правая ось была только у последней колонки;
    #   * шкала эффективности подстраивается под диапазон объекта, иначе при
    #     эффективности 1–2% линия ложится на ось и динамики не видно; чтобы
    #     разный масштаб не вводил в заблуждение, значение подписано числом
    #     в конце каждого месяца;
    #   * месяцы разделены вертикальной линией, «рисочкой» под осью и подложкой
    #     через месяц, названия месяцев — крупные и жирные;
    #   * третий (жёлтый) сегмент столбца — задачи «в работе»
    #     (NEW/WAITING/COMPLETING). Без него строки вроде коменданта, у которого
    #     все задачи месяца ещё открыты, выглядели как пустой график.
    @staticmethod
    def _month_groups(weeks):
        n = len(weeks)
        groups, cur, start = [], (weeks[0]["year"], weeks[0]["month"]), 0
        for i in range(1, n + 1):
            key = (weeks[i]["year"], weeks[i]["month"]) if i < n else None
            if key != cur:
                groups.append((start, i - 1, cur))
                if i < n:
                    cur, start = key, i
        return groups

    @staticmethod
    def _count_fmt(v, _=None):
        if v == 0:
            return "0"
        if v >= 1000:
            return ("%.0f тыс." % (v / 1000.0))
        return "%.0f" % v

    def grid_combo(self, weeks, entities, series_fn, name, ncols=3, nrows=3,
                   figsize=(12.33, 5.1), start_idx_fn=None):
        n = len(weeks)
        groups = self._month_groups(weeks)
        fig, axes = plt.subplots(nrows, ncols, figsize=figsize)
        axes = np.atleast_1d(axes).flatten()
        x = np.arange(n)
        for i, ax in enumerate(axes):
            if i >= len(entities):
                ax.axis("off")
                continue
            eid, ename = entities[i]
            completed, missed, completing, efficiency = series_fn(eid)
            si = start_idx_fn(eid) if start_idx_fn else 0
            c = np.array(completed, dtype=float)
            m = np.array(missed, dtype=float)
            w = np.array(completing, dtype=float)

            for gi, (s, e, (yy, mo)) in enumerate(groups):
                if gi % 2 == 1:
                    ax.axvspan(s - 0.5, e + 0.5, color=BG_LIGHT, zorder=0)
            ax.bar(x, c, color=GREEN, width=0.72, zorder=3)
            ax.bar(x, m, bottom=c, color=RED, width=0.72, zorder=3)
            ax.bar(x, w, bottom=c + m, color=AMBER, width=0.72, zorder=3)
            top = float(max((c + m + w).max(), 1))
            ax.set_ylim(0, top * 1.16)
            for sp in ("top", "right"):
                ax.spines[sp].set_visible(False)
            ax.spines["left"].set_color(GRID)
            ax.spines["bottom"].set_color("#C7D0E0")
            ax.spines["bottom"].set_linewidth(1.0)
            ax.set_axisbelow(True)
            ax.tick_params(axis="y", colors=TEXT_MUTED, labelsize=7.4, length=0, pad=2)
            ax.yaxis.set_major_locator(plt.MaxNLocator(3, integer=True))
            ax.yaxis.set_major_formatter(plt.FuncFormatter(self._count_fmt))
            ax.grid(axis="y", color=GRID, linewidth=0.7, zorder=1)
            ax.set_xlim(si - 0.6, n - 0.4)
            ax.set_xticks([])

            vals = [v for j, v in enumerate(efficiency) if v is not None and j >= si]
            emax = max(vals) if vals else 0.0
            span = max(emax * 1.45, 2.0)
            ax2 = ax.twinx()
            ax2.set_ylim(-span * 0.05, span)
            ax2.axhline(0, color=NAVY, alpha=0.10, linewidth=0.8, zorder=2)
            xs = [j for j, v in enumerate(efficiency) if v is not None and j >= si]
            ys = list(vals)
            if xs:
                ax2.plot(xs, ys, color=NAVY, linewidth=1.8, marker="o", markersize=2.6,
                         markerfacecolor="white", markeredgewidth=0.9, zorder=6)
            ticks = sorted({0.0, round(emax, 1)}) if emax > 0 else [0.0]
            ax2.set_yticks(ticks)
            ax2.set_yticklabels([ru("%g" % t) for t in ticks[:-1]] +
                                [ru("%g%%" % ticks[-1])])
            ax2.tick_params(axis="y", labelsize=7.0, colors=NAVY, length=0, pad=2)
            ax2.set_xlim(si - 0.6, n - 0.4)
            for sp in ax2.spines.values():
                sp.set_visible(False)

            for (s, e, (yy, mo)) in groups:
                if e < si:
                    continue
                if s > si:
                    ax.axvline(s - 0.5, color="#C7D0E0", linewidth=0.9, zorder=2)
                    ax.annotate("", xy=(s - 0.5, 0), xycoords=("data", "axes fraction"),
                                xytext=(0, -6), textcoords="offset points",
                                arrowprops=dict(arrowstyle="-", color="#8FA0BC", linewidth=1.0),
                                annotation_clip=False)
                ax.annotate(MONTHS_SHORT[mo].capitalize(),
                            xy=((max(s, si) + e) / 2.0, 0), xycoords=("data", "axes fraction"),
                            xytext=(0, -13), textcoords="offset points", ha="center", va="top",
                            fontsize=8.0, fontweight="bold", color=NAVY, annotation_clip=False)
                # Эффективность месяца подписываем ПОД осью, рядом с названием
                # месяца, а не на самой линии: у объектов с эффективностью
                # около нуля подписи на линии наезжали друг на друга.
                mc = sum(c[max(s, si):e + 1])
                mm_ = sum(m[max(s, si):e + 1])
                mv = (mc / (mc + mm_) * 100) if (mc + mm_) else None
                ax.annotate(ru("%.1f%%" % mv) if mv is not None else "—",
                            xy=((max(s, si) + e) / 2.0, 0), xycoords=("data", "axes fraction"),
                            xytext=(0, -22), textcoords="offset points", ha="center", va="top",
                            fontsize=6.8, color=NAVY, annotation_clip=False)

            tot_c, tot_m, tot_w = float(c[si:].sum()), float(m[si:].sum()), float(w[si:].sum())
            eff_total = round(tot_c / (tot_c + tot_m) * 100, 1) if (tot_c + tot_m) else None
            ax.set_title(textwrap.shorten(ename, width=42, placeholder="…"),
                         fontsize=8.6, color=TEXT_DARK, loc="left", fontweight="bold", pad=13)
            sub = "всего %s · вып. %s%s · эфф. %s" % (
                fmt_n(tot_c + tot_m + tot_w), fmt_n(tot_c),
                (" · в работе %s" % fmt_n(tot_w)) if tot_w else "",
                ("%.1f%%" % eff_total).replace(".", ",") if eff_total is not None else "—")
            ax.annotate(sub, xy=(0, 1.0), xycoords="axes fraction", xytext=(0, 3),
                        textcoords="offset points", fontsize=7.0, color=TEXT_MUTED, ha="left")
        fig.tight_layout(rect=[0, 0, 1, 1])
        fig.subplots_adjust(hspace=1.02, wspace=0.30)
        return self._save(fig, name, tight=False)

    # --- 5. дерево план/неплан (только фигуры, текст — нативный в PPTX) ---
    def tree_shapes(self, model, name="tree_shapes.png", figsize=(13.2, 6.0)):
        from matplotlib.patches import FancyBboxPatch
        tree = model["tree"]
        fig = plt.figure(figsize=figsize)
        ax = fig.add_axes([0, 0, 1, 1])
        ax.axis("off")
        ax.set_xlim(0, 100)
        ax.set_ylim(0, 100)

        def box(x, y, w, h, facecolor):
            ax.add_patch(FancyBboxPatch((x - w / 2, y - h / 2), w, h,
                                        boxstyle="round,pad=0.6,rounding_size=2.2",
                                        linewidth=0, facecolor=facecolor, zorder=2))

        root_x, root_y, root_w, root_h = 13, 50, 20, 22
        box(root_x, root_y, root_w, root_h, NAVY)
        cats = [
            {"y": 74, "data": tree["planned"], "color": BLUE,
             "eff": tree["planned_month_eff"], "due": tree["planned_month_due"]},
            {"y": 26, "data": tree["unplanned"], "color": AMBER,
             "eff": tree["unplanned_month_eff"], "due": tree["unplanned_month_due"]},
        ]
        box_x, box_w, box_h = 60, 62, 36
        box_left, box_right = box_x - box_w / 2, box_x + box_w / 2
        months = [m["label"] for m in tree["months_fold"]]
        for c in cats:
            ax.plot([root_x + root_w / 2, box_left], [50, c["y"]], color=TEXT_MUTED,
                    linewidth=1.4, zorder=1)
            box(box_x, c["y"], box_w, box_h, c["color"])
            spark_w, spark_h = box_w * 0.40, box_h * 0.52
            spark_left = box_right - spark_w - 3
            spark_bottom = c["y"] - spark_h / 2 + 1
            axin = ax.inset_axes([spark_left / 100, spark_bottom / 100,
                                  spark_w / 100, spark_h / 100])
            due_vals = c["due"]
            n = len(due_vals)
            axin.bar(range(n), due_vals, width=0.55, color=WHITE, alpha=0.28, zorder=1)
            # Столбцы занимают нижние ~38% высоты, линия эффективности — выше
            # (см. axin2 ниже). Иначе линия и её подписи ложатся ровно на
            # числа над столбцами и цифры не читаются.
            axin.set_ylim(0, max(due_vals) * 3.2 if due_vals and max(due_vals) > 0 else 1)
            axin.set_xlim(-0.35, n - 0.65)
            axin.set_xticks(range(len(months)))
            axin.set_xticklabels(months, fontsize=6.6, color=WHITE)
            axin.set_yticks([])
            axin.patch.set_alpha(0)
            for sp in axin.spines.values():
                sp.set_visible(False)
            axin.tick_params(length=0)
            for xi, v in enumerate(due_vals):
                axin.annotate(fmt_n(v), (xi, v), textcoords="offset points", xytext=(0, 3),
                              ha="center", va="bottom", fontsize=5.8, color=WHITE, alpha=0.85)
            axin2 = axin.twinx()
            vals = [v if v is not None else 0 for v in c["eff"]]
            axin2.plot(range(len(vals)), vals, color=WHITE, linewidth=2.3, marker="o",
                       markersize=4, markerfacecolor=WHITE, markeredgecolor=WHITE, zorder=4)
            # Линию кладём в ФИКСИРОВАННУЮ полосу высоты (55–92%), а не масштабируем
            # от нуля. Иначе при большом разбросе процентов (например 1 / 7 / 5 / 3)
            # низкие точки проваливаются вниз, ровно на подписи столбцов, и цифры
            # становятся нечитаемыми. Полоса не зависит от самих значений, поэтому
            # наложение не может возникнуть ни при каких данных.
            axin2.set_ylim(*band_ylim(vals, 0.55, 0.92))
            axin2.set_xlim(-0.35, len(vals) - 0.65)
            axin2.set_yticks([])
            axin2.patch.set_alpha(0)
            for sp in axin2.spines.values():
                sp.set_visible(False)
            axin2.tick_params(length=0)
            for xi, v in enumerate(vals):
                axin2.annotate("%.0f%%" % v, (xi, v), textcoords="offset points", xytext=(0, 6),
                               ha="center", fontsize=6.6, color=WHITE, fontweight="bold", zorder=5)
        return self._save(fig, name, tight=False)


# ================================================================== ТАБЛИЦЫ
def kom_name(koms, kid):
    k = koms.get(kid, {})
    nm = k.get("name") or kid
    nm = nm.capitalize() if isinstance(nm, str) else nm
    role = (k.get("roleName") or "").strip()
    return "%s (%s)" % (nm, role) if role else nm


def build_tables(model, keeps, kom_min=KOM_MIN_TASKS):
    """Таблицы приложений.

    Изменения прохода 17 (согласованы 26.08.2026):
      * в ячейке месяца три числа — вып/проп/В РАБОТЕ (было два). Без третьего
        строка, у которой все задачи месяца ещё открыты, выглядела как «0/0»;
      * у таблиц по задачам каждый месяц занимает ДВЕ колонки: объём
        (вып/проп/раб) и эффективность этого месяца — раньше эффективность была
        только итоговая за период, промежуточных значений видно не было;
      * состав строк берётся из ПОБЛОЧНЫХ списков активности (см. keeps).
    """
    months = model["meta"]["months"]
    month_keys = [m["key"] for m in months]
    month_labels = [m["label"] for m in months]
    objs = model["objects"]
    koms = model["komendanty"]
    first_activity = model.get("object_first_activity", {})
    tables = {}

    def month_cell(oid, mm, value_str):
        fa = first_activity.get(oid)
        if not fa:
            return value_str
        fad = dt.date.fromisoformat(fa)
        if (fad.year, fad.month) > (mm["year"], mm["month"]):
            return "—"
        if (fad.year, fad.month) == (mm["year"], mm["month"]):
            return "%s (с %s)" % (value_str, fad.strftime("%d.%m"))
        return value_str

    # --- обращения ---
    b2 = model["block2"]
    ids2 = [o for o in b2["objects_sorted"] if o in keeps["b2"]]
    rows = []
    for oid in ids2:
        mrow = b2["obj_month"].get(oid, {})
        cells = [month_cell(oid, mm, fmt_n(mrow.get(mm["key"], 0))) for mm in months]
        rows.append([objs.get(oid, oid)] + cells + [fmt_n(b2["obj_total"][oid])])
    tables["b2"] = {
        "header": ["Объект"] + month_labels + ["Итого"],
        "rows": rows,
        "total_row": ["Итого"] + [fmt_n(sum(b2["obj_month"].get(o, {}).get(mk, 0) for o in ids2))
                                  for mk in month_keys]
                     + [fmt_n(sum(b2["obj_total"][o] for o in ids2))],
    }

    def cell3(mv):
        return "%s/%s/%s" % (fmt_n(mv.get("completed", 0)), fmt_n(mv.get("missed", 0)),
                             fmt_n(mv.get("completing", 0)))

    def cell_eff(mv):
        c, m = mv.get("completed", 0), mv.get("missed", 0)
        return ru("%.1f%%" % (c / (c + m) * 100)) if (c + m) else "—"

    ZERO = {"completed": 0, "missed": 0, "completing": 0}

    def month_table(block, ids, name_fn, key_field, tot_field,
                    use_first_activity=False, split_eff=False):
        out = []
        for eid in ids:
            mrow = block[key_field].get(eid, {})
            cells = []
            for mm in months:
                mv = mrow.get(mm["key"], ZERO)
                val = cell3(mv)
                cells.append(month_cell(eid, mm, val) if use_first_activity else val)
                if split_eff:
                    cells.append(cell_eff(mv))
            out.append([name_fn(eid)] + cells
                       + [ru("%.1f%%" % block[tot_field][eid]["efficiency"])])
        return out

    obj_name = lambda x: objs.get(x, x)
    kname = lambda x: kom_name(koms, x)
    hdr_kom = ["Комендант"] + ["%s\n(вып/проп/раб)" % m for m in month_labels] + ["Эфф. за период"]
    hdr_obj_split = ["Объект"]
    for m in month_labels:
        hdr_obj_split += ["%s\nвып/проп/раб" % m, "%s\nэфф." % m]
    hdr_obj_split += ["Эфф. за период"]
    hdr_obj = ["Объект"] + ["%s\n(вып/проп/раб)" % m for m in month_labels] + ["Эфф. за период"]

    for tkey, block, ids, name_fn, hdr, kf, tf, fa, split in (
        ("b3", model["block3"], [o for o in model["block3"]["objects_sorted"] if o in keeps["b3"]],
         obj_name, hdr_obj_split, "obj_month", "obj_total", True, True),
        ("b3p", model.get("block3p", model["block3"]),
         [o for o in model.get("block3p", model["block3"])["objects_sorted"] if o in keeps["b3p"]],
         obj_name, hdr_obj_split, "obj_month", "obj_total", True, True),
        ("b3u", model.get("block3u", model["block3"]),
         [o for o in model.get("block3u", model["block3"])["objects_sorted"] if o in keeps["b3u"]],
         obj_name, hdr_obj_split, "obj_month", "obj_total", True, True),
        ("b4", model["block4"],
         [k for k in model["block4"]["komendanty_sorted"]
          if k in keeps["b4"] and model["block4"]["kom_total"][k]["due"] >= kom_min],
         kname, hdr_kom, "kom_month", "kom_total", False, False),
        ("b4b", model["block4b"],
         [k for k in model["block4b"]["komendanty_sorted"]
          if k in keeps["b4b"] and model["block4b"]["kom_total"][k]["due"] >= kom_min],
         kname, hdr_kom, "kom_month", "kom_total", False, False),
        ("b5", model["block5"], [o for o in model["block5"]["objects_sorted"] if o in keeps["b5"]],
         obj_name, hdr_obj, "obj_month", "obj_total", True, False),
        ("b5b", model["block5b"], [o for o in model["block5b"]["objects_sorted"] if o in keeps["b5b"]],
         obj_name, hdr_obj, "obj_month", "obj_total", True, False),
    ):
        if not ids:
            continue
        tables[tkey] = {"header": hdr,
                        "rows": month_table(block, ids, name_fn, kf, tf,
                                            use_first_activity=fa, split_eff=split)}
    return tables


# ================================================================== ВЁРСТКА
class Deck:
    def __init__(self, period_str, period_short):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W_IN)
        self.prs.slide_height = Inches(SLIDE_H_IN)
        self.blank = self.prs.slide_layouts[6]
        self.period_str = period_str
        self.period_short = period_short
        self.page = 0

    # --- примитивы ---
    @staticmethod
    def _c(hexstr):
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    def slide(self):
        return self.prs.slides.add_slide(self.blank)

    def rect(self, s, x, y, w, h, fill=None, shape=MSO_SHAPE.RECTANGLE, radius=None):
        sh = s.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        if radius is not None:
            try:
                sh.adjustments[0] = radius
            except Exception:
                pass
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = self._c(fill)
        else:
            sh.fill.background()
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def round_rect(self, s, x, y, w, h, fill, radius=0.08):
        return self.rect(s, x, y, w, h, fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=radius)

    def text(self, s, x, y, w, h, text, size=14, bold=False, colorhex=TEXT_DARK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0, italic=False, wrap=True):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        for i, line in enumerate(str(text).split("\n")):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = line_spacing
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.italic = italic
            r.font.name = FONT
            r.font.color.rgb = self._c(colorhex)
        return tb

    def badge(self, s, x, y, num, d=0.5, fontsize=18):
        sh = self.rect(s, x, y, d, d, NAVY, shape=MSO_SHAPE.OVAL)
        tf = sh.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(num)
        r.font.size = Pt(fontsize)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = self._c(WHITE)
        return sh

    def header(self, s, num, title, subtitle):
        # Ширина под заголовок — вся строка до правого поля (было 9,5"), кегль
        # подбирается под длину: длинные заголовки слайдов 6/7/8 переносились
        # на вторую строку и накрывали подзаголовок.
        self.badge(s, MARGIN_IN, 0.42, num)
        w = SLIDE_W_IN - MARGIN_IN - (MARGIN_IN + 0.68)
        self.text(s, MARGIN_IN + 0.68, 0.35, w, 0.5, title,
                  size=fit_size(title, w, 22, 14.5, bold=True), bold=True, colorhex=NAVY, wrap=False)
        self.text(s, MARGIN_IN + 0.68, 0.82, w, 0.4, subtitle,
                  size=fit_size(subtitle, w, 12.5, 9.5), colorhex=TEXT_MUTED, wrap=False)

    def footer(self, s, dark=False):
        self.page += 1
        col = TEXT_ON_NAVY_MUTED if dark else TEXT_MUTED
        self.text(s, MARGIN_IN, SLIDE_H_IN - 0.38, 8, 0.3,
                  "Bruno · ВТБ · %s" % self.period_short, size=9, colorhex=col)
        self.text(s, SLIDE_W_IN - MARGIN_IN - 1, SLIDE_H_IN - 0.38, 1, 0.3, str(self.page),
                  size=9, colorhex=col, align=PP_ALIGN.RIGHT)

    def kpi(self, s, x, y, w, h, value, label, accent=NAVY, size=26):
        # Значение и подпись — строго в одну строку с подбором кегля: пара
        # «всего/выполнено» вида «47 122/ 8 308» при 26 pt не влезала в плашку
        # и второй строкой ложилась на подпись.
        self.round_rect(s, x, y, w, h, BG_LIGHT, radius=0.10)
        vw, lw = w - 0.3, w - 0.3
        self.text(s, x + 0.15, y + 0.12, vw, h - 0.55, value,
                  size=fit_size(value, vw, size, 12.0, bold=True), bold=True,
                  colorhex=accent, wrap=False)
        self.text(s, x + 0.15, y + h - 0.42, lw, 0.36, label,
                  size=fit_size(label, lw, 10.5, 7.6), colorhex=TEXT_MUTED,
                  line_spacing=0.95, wrap=False)

    def picture(self, s, path, x, y, w, h):
        """Вписывает картинку в бокс с сохранением пропорций; возвращает
        фактический прямоугольник — по нему позиционируются нативные подписи."""
        from PIL import Image
        iw, ih = Image.open(path).size
        ar_img, ar_box = iw / ih, w / h
        if ar_img > ar_box:
            nw, nh = w, w / ar_img
        else:
            nh, nw = h, h * ar_img
        nx, ny = x + (w - nw) / 2, y + (h - nh) / 2
        s.shapes.add_picture(path, Inches(nx), Inches(ny), Inches(nw), Inches(nh))
        return nx, ny, nw, nh

    # --- легенда / подписи осей ---
    def legend_row(self, s, x, y, w, items, fontsize=8.7, sq=0.13, gap_after_sq=0.08,
                   item_gap=0.28, row_h=0.24, align=PP_ALIGN.CENTER):
        def text_w(t):
            return 0.072 * fontsize / 9.0 * len(t) + sq + gap_after_sq + item_gap

        rows, cur, cur_w = [], [], 0.0
        for hexcolor, label in items:
            iw = text_w(label)
            if cur and cur_w + iw > w:
                rows.append(cur)
                cur, cur_w = [], 0.0
            cur.append((hexcolor, label, iw))
            cur_w += iw
        if cur:
            rows.append(cur)
        yy = y
        for row in rows:
            row_w = sum(it[2] for it in row) - item_gap
            xx = x + (w - row_w) / 2 if align == PP_ALIGN.CENTER else x
            for hexcolor, label, iw in row:
                self.rect(s, xx, yy + (row_h - sq) / 2, sq, sq, hexcolor)
                self.text(s, xx + sq + gap_after_sq, yy, iw, row_h, label, size=fontsize,
                          colorhex=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)
                xx += iw
            yy += row_h
        return yy - y

    def axis_caption(self, s, x, y, w, txt, align=PP_ALIGN.LEFT, colorhex=TEXT_MUTED, size=8.7):
        self.text(s, x, y, w, 0.22, txt, size=size, colorhex=colorhex, align=align)

    def chart_legend(self, s, rect, legend_items, y, line_label="Эффективность, %"):
        """Легенда под графиком: квадратики по центру + отрезок линии
        эффективности справа. Геометрия — как в утверждённой колоде."""
        nx, ny, nw, nh = rect
        self.legend_row(s, nx, y, nw * 0.72, legend_items)
        ln = s.shapes.add_connector(1, Inches(nx + nw * 0.76), Inches(y + 0.11),
                                    Inches(nx + nw * 0.76 + 0.16), Inches(y + 0.11))
        ln.line.color.rgb = self._c(NAVY)
        ln.line.width = Pt(2.2)
        self.text(s, nx + nw * 0.76 + 0.22, y, nw * 0.24, 0.24, line_label, size=8.69,
                  colorhex=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    def combo_chrome(self, s, rect, ylabel_left, ylabel_right, legend_items,
                     line_label="Эффективность, %"):
        nx, ny, nw, nh = rect
        self.axis_caption(s, nx, ny - 0.20, nw / 2, ylabel_left)
        self.axis_caption(s, nx + nw / 2, ny - 0.20, nw / 2, ylabel_right,
                          align=PP_ALIGN.RIGHT, colorhex=NAVY)
        yy = ny + nh + 0.10
        self.legend_row(s, nx, yy, nw * 0.72, legend_items)
        ln = s.shapes.add_connector(1, Inches(nx + nw * 0.76), Inches(yy + 0.11),
                                    Inches(nx + nw * 0.76 + 0.16), Inches(yy + 0.11))
        ln.line.color.rgb = self._c(NAVY)
        ln.line.width = Pt(2.2)
        self.text(s, nx + nw * 0.76 + 0.22, yy, nw * 0.24, 0.24, line_label, size=8.7,
                  colorhex=TEXT_DARK, anchor=MSO_ANCHOR.MIDDLE)

    # --- таблицы ---
    def table(self, s, x, y, w, h, header, rows, col0_ratio=0.24, font_size=9,
              col_weights=None):
        """col_weights — относительные ширины колонок (кроме первой). Нужны
        таблицам «месяц = объём + эффективность»: колонка с тремя числами
        («4 829/91 986/1 498») втрое шире колонки с процентом, а при равной
        ширине она переносилась на две строки."""
        shape = s.shapes.add_table(len(rows) + 1, len(header), Inches(x), Inches(y),
                                   Inches(w), Inches(h))
        tbl = shape.table
        col0 = w * col0_ratio
        rest = w - col0
        n_rest = max(len(header) - 1, 1)
        if col_weights and len(col_weights) == n_rest:
            tot = float(sum(col_weights))
            widths = [rest * cw / tot for cw in col_weights]
        else:
            widths = [rest / n_rest] * n_rest
        tbl.columns[0].width = Inches(col0)
        for c in range(1, len(header)):
            tbl.columns[c].width = Inches(widths[c - 1])
        for c, htext in enumerate(header):
            tbl.cell(0, c).text = htext
        for r, row in enumerate(rows):
            for c, val in enumerate(row):
                tbl.cell(r + 1, c).text = str(val)
        for ci, cell in enumerate(tbl.rows[0].cells):
            cell.fill.solid()
            cell.fill.fore_color.rgb = self._c(NAVY)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.06)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.bold = True
                    r.font.size = Pt(font_size)
                    r.font.name = FONT
                    r.font.color.rgb = self._c(WHITE)
        for ri in range(1, len(tbl.rows)):
            for ci, cell in enumerate(tbl.rows[ri].cells):
                cell.fill.solid()
                cell.fill.fore_color.rgb = self._c(BG_LIGHT if ri % 2 == 0 else WHITE)
                cell.vertical_anchor = MSO_ANCHOR.MIDDLE
                cell.margin_left = cell.margin_right = Inches(0.06)
                for p in cell.text_frame.paragraphs:
                    p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
                    for r in p.runs:
                        r.font.size = Pt(font_size)
                        r.font.name = FONT
                        r.font.color.rgb = self._c(TEXT_DARK)
                        r.font.bold = (ci == 0)
        return tbl


# ================================================================== СБОРКА
STD_LEGEND = [(GREEN, "Выполнено"), (RED, "Пропущено")]
STD_LEGEND_NC = [(GREEN, "Выполнено"), (RED, "Не выполнено")]
GRID_LEGEND = [(GREEN, "Выполнено"), (RED, "Пропущено"), (AMBER, "В работе"),
               (NAVY, "Эффективность, % (правая шкала своя у каждого графика)")]
GRID_LEGEND_NC = GRID_LEGEND


def combo_legend(has_completing, missed_label="Пропущено"):
    items = [(GREEN, "Выполнено"), (RED, missed_label)]
    if has_completing:
        items.append((AMBER, "В работе"))
    return items


def weekly_eff(week):
    return [round(c / (c + m) * 100, 1) if (c + m) else None
            for c, m in zip(week["completed"], week["missed"])]


# Сетка приложений — решение заказчика 26.08.2026: 3 колонки, максимум 9
# графиков на слайд, остальные переносятся на следующий слайд. Раньше было до
# 16 на слайде, отчего всё и было нечитаемым.
GRID_COLS = 3
GRID_ROWS = 3
GRID_PAGE = GRID_COLS * GRID_ROWS
# Размер ОДНОЙ ячейки в дюймах. Совпадает с местом на слайде (см. GRID_BOX):
# фигура делается ровно того размера, в котором ляжет на слайд, поэтому кегли
# matplotlib = кегли на слайде, ничего не ужимается.
CELL_W, CELL_H = 4.11, 1.78
GRID_BOX = (MARGIN_IN, 1.40, CELL_W * GRID_COLS, CELL_H * GRID_ROWS)


def grid_pages(entities):
    return [entities[i:i + GRID_PAGE] for i in range(0, len(entities), GRID_PAGE)] or []


def grid_dims(n):
    ncols = min(GRID_COLS, max(n, 1))
    return ncols, max(math.ceil(n / ncols), 1)


def grid_figsize(ncols, nrows):
    return (CELL_W * ncols, CELL_H * nrows)




# ------------------------------------------------------------------ геометрия
# Все координаты сняты с утверждённой колоды (БРУНО_исправленно.pptx) шейп в
# шейп. Менять их без согласования нельзя: заказчик сверяет копию с оригиналом.
HALF_W = 6.02
LX, RX = 0.50, 6.82


def build(model, out_path, imgdir):
    C = Charts(imgdir)
    weeks = model["weeks"]
    objs = model["objects"]
    koms = model["komendanty"]
    meta = model["meta"]

    d_from = dt.date.fromisoformat(meta["period_from"])
    d_to = dt.date.fromisoformat(meta["period_to"])
    period_str = "%s — %s" % (d_from.strftime("%d.%m.%Y"), d_to.strftime("%d.%m.%Y"))
    period_short = "%s–%s" % (d_from.strftime("%d.%m"), d_to.strftime("%d.%m.%Y"))
    D = Deck(period_str, period_short)

    keep_objects = set(model["active_objects_current_month"])
    keep_koms = set(model["active_komendanty_current_month"])

    # Поблочная активность в текущем месяце (проход 17). Раньше во всех
    # приложениях применялся один общий список, и в приложение по неплановым
    # задачам попадали строки без единой неплановой задачи в текущем месяце.
    ACT = model.get("active_current_month") or {}

    def keep(block_key, fallback):
        s = ACT.get(block_key)
        return set(s) if s is not None else set(fallback)

    b2, b3 = model["block2"], model["block3"]
    b4, b4b = model["block4"], model["block4b"]
    b5, b5b = model["block5"], model["block5b"]
    b6 = model["block6"]
    tree = model["tree"]

    week_starts = [dt.date.fromisoformat(w["start"]) for w in weeks]

    def start_idx(oid):
        fa = model["object_first_activity"].get(oid)
        if not fa:
            return 0
        d = dt.date.fromisoformat(fa)
        for i in range(len(week_starts) - 1, -1, -1):
            if week_starts[i] <= d:
                return i
        return 0

    # Ряды для мини-графиков: выполнено / пропущено / В РАБОТЕ / эффективность.
    # Третий ряд добавлен 26.08.2026: задачи в статусах NEW/WAITING/COMPLETING
    # раньше нигде не рисовались, и строка, у которой все задачи месяца ещё
    # открыты (напр. комендант o.duka: 0/0/299 за август), выглядела пустой.
    def obj_series(block):
        def fn(oid):
            w = block["obj_week"][oid]
            return w["completed"], w["missed"], w["completing"], weekly_eff(w)
        return fn

    def kom_series(block):
        def fn(kid):
            w = block["kom_week"][kid]
            return w["completed"], w["missed"], w["completing"], weekly_eff(w)
        return fn

    def eff(total):
        """Единственная формула эффективности во всей колоде (решение заказчика
        22.08.2026): выполнено / (выполнено + пропущено). В утверждённой версии
        на слайдах 5 и 6 под словом «выполнено» стояла другая формула — доля от
        всех заведённых задач; здесь она больше не используется."""
        return total["efficiency"]

    # ------------------------------------------------------------- картинки
    img = {}
    img["tree"] = C.tree_shapes(model)
    img["b2_main"] = C.feedback_area(model)
    img["b3_main"], b3_has_cmp = C.combo(
        weeks, b3["total_week"]["completed"], b3["total_week"]["missed"],
        b3["total_week"]["completing"], b3["total_efficiency_week"], "b3_main.png")
    tw4b = b4b["total_week"]
    img["b4b_main"], b4b_has_cmp = C.combo(
        weeks, tw4b["completed"], tw4b["missed"], tw4b["completing"],
        b4b["total_efficiency_week"], "b4b_main.png")
    tw4 = b4["total_week"]
    # Было two_part: «в работе» подмешивалось в красное «не выполнено». Теперь
    # третий сегмент показан отдельно — как во всех остальных графиках колоды
    # (решение 26.08.2026).
    img["b4_main"], b4_has_cmp = C.combo(
        weeks, tw4["completed"], tw4["missed"], tw4["completing"],
        b4["total_efficiency_week"], "b4_main.png")
    tw5 = b5["total_week"]
    img["b5_main"], b5_has_cmp = C.combo(
        weeks, tw5["completed"], tw5["missed"], tw5["completing"],
        b5["total_efficiency_week"], "b5_main.png")
    tw5b = b5b["total_week"]
    img["b5b_main"], b5b_has_cmp = C.combo(
        weeks, tw5b["completed"], tw5b["missed"], tw5b["completing"],
        b5b["total_efficiency_week"], "b5b_main.png")

    keep_b3 = keep("b3", keep_objects)
    keep_b3p = keep("b3p", keep_objects)
    keep_b3u = keep("b3u", keep_objects)
    keep_b5 = keep("b5", keep_objects)
    keep_b5b = keep("b5b", keep_objects)
    keep_k4 = keep("b4", keep_koms)
    keep_k4b = keep("b4b", keep_koms)

    ent3 = [(o, objs.get(o, o)) for o in b3["objects_sorted"] if o in keep_b3]
    ent4 = [(k, kom_name(koms, k)) for k in b4["komendanty_sorted"]
            if k in keep_k4 and b4["kom_total"][k]["due"] >= KOM_MIN_TASKS]
    ent4b = [(k, kom_name(koms, k)) for k in b4b["komendanty_sorted"]
             if k in keep_k4b and b4b["kom_total"][k]["due"] >= KOM_MIN_TASKS]
    ent5 = [(o, objs.get(o, o)) for o in b5["objects_sorted"] if o in keep_b5]
    ent5b = [(o, objs.get(o, o)) for o in b5b["objects_sorted"] if o in keep_b5b]
    # Каждый набор режется на страницы по 9 графиков (3x3) — img[key] хранит
    # СПИСОК путей к картинкам, по одному на слайд.
    for key, ents, series, sfn in (
        ("b3_grid", ent3, obj_series(b3), start_idx),
        ("b4_grid", ent4, kom_series(b4), None),
        ("b4b_grid", ent4b, kom_series(b4b), None),
        ("b5_grid", ent5, obj_series(b5), start_idx),
        ("b5b_grid", ent5b, obj_series(b5b), start_idx),
    ):
        if not ents:
            continue
        pages = []
        for pi, chunk in enumerate(grid_pages(ents)):
            nc, nr = grid_dims(len(chunk))
            pages.append((C.grid_combo(weeks, chunk, series, "%s_%d.png" % (key, pi),
                                       ncols=nc, nrows=nr, figsize=grid_figsize(nc, nr),
                                       start_idx_fn=sfn), nr))
        img[key] = pages

    tables = build_tables(model, {"b2": keep("b2", keep_objects), "b3": keep_b3,
                                  "b3p": keep_b3p, "b3u": keep_b3u,
                                  "b4": keep_k4, "b4b": keep_k4b,
                                  "b5": keep_b5, "b5b": keep_b5b})

    # В утверждённой колоде здесь стоял жёсткий список объектов «со статусом
    # Готово»; по согласованию 21.08.2026 фильтр стал динамическим, поэтому
    # сноска описывает правило, а не перечисляет объекты — иначе она разрастается
    # на две строки и наезжает на футер.

    # =================================================== 1. Титул
    s = D.slide()
    D.rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, NAVY)
    D.rect(s, 0, 0, SLIDE_W_IN, 0.06, BLUE)
    D.text(s, 1.00, 2.55, 11.30, 1.00, "BRUNO", size=54, bold=True, colorhex=WHITE)
    D.text(s, 1.00, 3.55, 11.30, 0.60, "Отчёт по эксплуатации объектов", size=22, colorhex=WHITE)
    D.text(s, 1.00, 4.15, 11.30, 0.50,
           "Оклейка QR · обращения · плановые и неплановые задачи · коменданты · уборка",
           size=13, colorhex=TEXT_ON_NAVY_MUTED)
    D.text(s, 1.00, 6.55, 11.30, 0.40, "ВТБ  ·  %s" % period_str, size=13, bold=True, colorhex=BLUE)
    D.page += 1

    n_obj_active = len(set(b2["obj_total"]) | set(b3["obj_total"]) | set(b5["obj_total"]))
    n_obj_dir = len(objs)

    # =================================================== 2. Обзор отчёта
    s = D.slide()
    D.text(s, MARGIN_IN, 0.45, 11.00, 0.60, "Обзор отчёта", size=26, bold=True, colorhex=NAVY)
    D.text(s, MARGIN_IN, 1.05, 11.00, 0.40,
           "%d активных объектов портфеля · период %s · данные Bruno API" % (n_obj_active, period_str),
           size=12.5, colorhex=TEXT_MUTED)
    overview = [
        ("1", "Внедрение", "Ход оклейки QR-кодов по зданиям портфеля"),
        ("2", "Плановые / неплановые", "Структура портфеля задач и конверсия обращений"),
        ("3", "Обращения", "Динамика входящих обращений сотрудников по QR, по объектам"),
        ("4", "Задачи", "Плановые и неплановые задачи, эффективность выполнения"),
        ("5", "Коменданты", "Плановые и неплановые задачи — динамика по каждому коменданту"),
        ("6", "Уборка", "Плановые и неплановые задачи по объектам, эффективность выполнения"),
    ]
    for i, (num, title, desc) in enumerate(overview):
        x = 0.50 + (i % 3) * 4.00
        y = 1.70 + (i // 3) * 2.60
        D.round_rect(s, x, y, 3.78, 2.38, BG_LIGHT, radius=0.06)
        D.badge(s, x + 0.22, y + 0.20, num, d=0.44, fontsize=15)
        D.text(s, x + 0.22, y + 0.80, 3.34, 0.40, title, size=15.5, bold=True, colorhex=NAVY)
        D.text(s, x + 0.22, y + 1.20, 3.34, 1.06, desc, size=10.8, colorhex=TEXT_DARK, line_spacing=1.05)
    D.text(s, MARGIN_IN, 6.98, 12.33, 0.30,
           "Bruno: в справочнике %d объектов, %d из них без активности за период (тест/дубли) — "
           "исключены из отчёта. Счётчики объектов на слайдах различаются: каждый считает своё "
           "множество." % (n_obj_dir, max(n_obj_dir - n_obj_active, 0)),
           size=8.6, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 3. Внедрение (ручные данные)
    s = D.slide()
    D.header(s, 1, "Внедрение: оклейка QR-кодов", "Статус по зданиям портфеля на дату отчёта")
    for i, (title, items, accent, deadline) in enumerate(IMPLEMENTATION):
        x = 0.50 + i * 4.02
        D.round_rect(s, x, 1.75, 3.78, 4.85, BG_LIGHT, radius=0.05)
        D.rect(s, x + 0.22, 1.97, 0.10, 0.62, accent)
        D.text(s, x + 0.44, 2.06, 3.18, 0.40, title, size=15, bold=True, colorhex=TEXT_DARK)
        D.text(s, x + 3.09, 2.01, 0.60, 0.50, str(len(items)), size=30, bold=True, colorhex=NAVY)
        if deadline:
            D.text(s, x + 0.46, 2.35, 3.18, 0.35, deadline, size=10.5, italic=True, colorhex=TEXT_MUTED)
        yy = 3.00
        for it in items:
            D.text(s, x + 0.30, yy, 3.18, 0.42, "•  " + it, size=11, colorhex=TEXT_DARK, line_spacing=1.05)
            yy += 0.50
    D.text(s, MARGIN_IN, 6.75, 11.30, 0.16, IMPLEMENTATION_NOTE,
           size=9.5, bold=True, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 4. Дерево план/неплан
    s = D.slide()
    D.header(s, 2, "Плановые и неплановые задачи",
             "Структура портфеля по типу постановки и конверсия обращений")
    rect = D.picture(s, img["tree"], 1.33, 1.55, 10.67, 4.85)
    tree_text(D, s, rect, model)
    D.text(s, MARGIN_IN, 6.88, 11.30, 0.14,
           "Конверсия обращений в задачи может превышать 100% — необработанное обращение может "
           "формироваться повторно на следующий день выгрузки.",
           size=8.5, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 5. Обращения
    s = D.slide()
    D.header(s, 3, "Динамика обращений",
             "Обращения сотрудников по QR-кодам, по объектам, по неделям")
    for i, (val, lbl) in enumerate(((fmt_n(b2["total"]), "обращений всего"),
                                    (fmt_n(b2["max_week"]), "макс. в неделю"),
                                    (str(b2["n_objects"]), "объектов с обращениями"))):
        D.kpi(s, 0.50 + i * 2.75, 1.70, 2.55, 1.15, val, lbl)
    rect = D.picture(s, img["b2_main"], 1.62, 3.24, 9.62, 2.65)
    D.axis_caption(s, 1.62, 3.04, 9.62, "Обращений в неделю")
    legend_items = [(cat_color(i), objs.get(oid, oid)) for i, oid in enumerate(b2["objects_sorted"])]
    D.legend_row(s, MARGIN_IN, 6.02, SLIDE_W_IN - 2 * MARGIN_IN, legend_items,
                 fontsize=7.6, row_h=0.22)
    D.text(s, MARGIN_IN, 6.78, 11.30, 0.30,
           "Резкий рост числа обращений сигнализирует о завершении оклейки QR в здании, "
           "а не об ухудшении качества обслуживания.", size=8.6, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 6. Все задачи
    s = D.slide()
    D.header(s, 4, "Плановые и неплановые задачи: объём и эффективность",
             "Объём портфеля задач и эффективность фактического выполнения")
    tot3 = b3["total"]
    D.kpi(s, 2.06, 1.68, 2.55, 1.15, short_num(tot3["due"]), "задач всего")
    D.kpi(s, 4.81, 1.68, 2.55, 1.15, ru("%.1f%%" % tot3["efficiency"]), "эффективность за период",
          accent=RED if tot3["efficiency"] < 30 else GREEN)
    D.kpi(s, 7.56, 1.68, 2.55, 1.15, fmt_n(tot3["completed"]), "выполненных задач")
    rect = D.picture(s, img["b3_main"], 1.68, 3.33, 8.94, 2.35)
    D.axis_caption(s, 1.68, 3.13, 4.47, "Задач в неделю")
    D.axis_caption(s, 6.15, 3.13, 4.47, "Эффективность, %", align=PP_ALIGN.RIGHT, colorhex=NAVY)
    D.chart_legend(s, rect, combo_legend(b3_has_cmp), y=5.83)
    D.text(s, MARGIN_IN, 6.86, 11.30, 0.35,
           "Картина сильно различается по объектам — см. детализацию в приложении. "
           "Эффективность = выполнено / (выполнено + пропущено).",
           size=8.6, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 7. Коменданты: план + неплан
    s = D.slide()
    D.header(s, 5, "Коменданты: плановые и неплановые задачи",
             "Слева — плановые задачи, отработанные комендантами лично. "
             "Справа — неплановые (обращения/заявки)")
    for x, title, block, imgkey, legend in (
        (LX, "Коменданты · плановые", b4b, "b4b_main", combo_legend(b4b_has_cmp)),
        (RX, "Коменданты · неплановые", b4, "b4_main", combo_legend(b4_has_cmp)),
    ):
        t = block["total"]
        D.text(s, x, 1.83, HALF_W, 0.30, title, size=13.5, bold=True, colorhex=NAVY)
        D.kpi(s, x, 2.18, 2.95, 0.95, "%s/ %s" % (fmt_n(t["due"]), fmt_n(t["completed"])),
              "задач всего/ выполнено")
        D.kpi(s, x + 3.07, 2.18, 2.95, 0.95, ru("%.1f%%" % eff(t)), "эффективность")
        D.axis_caption(s, x, 3.41, 3.01, "Задач/нед.")
        D.axis_caption(s, x + 3.01, 3.41, 3.01, "Эфф., %", align=PP_ALIGN.RIGHT, colorhex=NAVY)
        r = D.picture(s, img[imgkey], x, 3.61, HALF_W, 1.59)
        D.chart_legend(s, r, legend, y=5.30)
    D.text(s, MARGIN_IN, 6.86, 11.30, 0.35, "Детализация по каждому коменданту — в приложении.",
           size=8.6, italic=True, colorhex=TEXT_MUTED)
    D.footer(s)

    # =================================================== 8. Уборка: план + неплан + скорость
    # Единственное согласованное отступление от утверждённой колоды: справа
    # добавлен график НЕПЛАНОВЫХ задач по уборке, под каждым графиком — своя
    # таблица скорости (задач на человека и ФАКТИЧЕСКИХ минут на задачу).
    s = D.slide()
    D.header(s, 6, "Уборка: плановые и неплановые задачи",
             "Плановые и неплановые задачи по объектам + скорость по башням (день/вечер)")
    for x, title, block, imgkey, legend, kind in (
        (LX, "Уборка · плановые", b5, "b5_main", combo_legend(b5_has_cmp), "planned"),
        (RX, "Уборка · неплановые", b5b, "b5b_main", combo_legend(b5b_has_cmp), "unplanned"),
    ):
        t = block["total"]
        share = eff(t)
        D.text(s, x, 1.55, HALF_W, 0.30, title, size=13.5, bold=True, colorhex=NAVY)
        D.kpi(s, x, 1.88, 2.95, 0.95, "%s/ %s" % (fmt_n(t["due"]), fmt_n(t["completed"])),
              "задач всего/ выполнено")
        D.kpi(s, x + 3.07, 1.88, 2.95, 0.95, ru("%.1f%%" % share), "эффективность",
              accent=RED if share < 30 else GREEN)
        D.axis_caption(s, x, 3.08, 3.01, "Задач/нед.")
        D.axis_caption(s, x + 3.01, 3.08, 3.01, "Эфф., %", align=PP_ALIGN.RIGHT, colorhex=NAVY)
        r = D.picture(s, img[imgkey], x, 3.28, HALF_W, 1.42)
        D.chart_legend(s, r, legend, y=4.80)

        avg = b6["average"][kind]
        n_tow = b6["n_towers_with_any_data"][kind]
        D.text(s, x, 5.12, HALF_W, 0.24,
               "Скорость работы по 3 башням (с данными: %d из %d)" % (n_tow, b6["n_towers_total"]),
               size=11, bold=True, colorhex=NAVY)
        # Колонки прохода 17 (решение заказчика 26.08.2026):
        #   «Назначено/чел.» — ВСЕ задачи смены на одного вышедшего исполнителя;
        #   «Выполнено/чел.» — то, что было в колоде раньше под именем
        #   «Задач/чел.» (только завершённые задачи), из-за чего показатель и
        #   выглядел заниженным в 10-20 раз;
        #   «мин./задачу» — среднее взвешенное И медиана: среднее вытягивают
        #   единичные «висящие» задачи, медиана показывает типичный случай.
        rows = []
        for shift, shift_label in (("day", "День"), ("evening", "Вечер")):
            a = avg[shift]
            rows.append([shift_label,
                         num_cell(a.get("assigned_per_employee")),
                         num_cell(a["avg_tasks_per_employee"]),
                         num_cell(a["avg_duration_min"]),
                         num_cell(a.get("median_duration_min"))])
        D.table(s, x, 5.42, HALF_W, 1.05,
                ["Смена", "Назначено\nна чел.", "Выполнено\nна чел.",
                 "Мин./зад.\nсреднее", "Мин./зад.\nмедиана"], rows,
                col0_ratio=0.20, font_size=7.6)
    D.text(s, MARGIN_IN, 6.62, 12.33, 0.50,
           "«Назначено/выполнено на человека» — задачи смены, делённые на число исполнителей, "
           "которые в эту смену закрыли хотя бы одну задачу; смена определяется по плановому "
           "времени задачи.\n«Мин./задачу» — фактическое время от «начал» до «завершил»: среднее "
           "взвешено по числу задач, медиана устойчива к единичным «висящим» задачам. "
           "Плановые — роль уборщиц; неплановые — вместе с «Менеджером клининга».",
           size=8.0, italic=True, colorhex=TEXT_MUTED, line_spacing=1.12)
    D.footer(s)

    # =================================================== 9. Итоги периода
    s = D.slide()
    D.rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, NAVY)
    D.text(s, MARGIN_IN, 0.50, 11.00, 0.60, "Итоги периода", size=28, bold=True, colorhex=WHITE)
    D.text(s, MARGIN_IN, 1.15, 11.00, 0.40, period_str, size=14, colorhex=TEXT_ON_NAVY_MUTED)
    summary = [
        ("%d из %d" % (len(IMPLEMENTATION[0][1]), IMPLEMENTATION_TOTAL), "зданий полностью оклеены QR"),
        (fmt_n(b2["total"]), "обращений принято через QR"),
        (short_num(b3["total"]["due"]), "задач всего"),
        (ru("%.1f%%/ %.1f%%" % (tree["planned"]["efficiency"], tree["unplanned"]["efficiency"])),
         "эффективность по плановым / неплановым  задачам"),
        (ru("%.1f%%/ %.1f%%" % (eff(b4b["total"]), eff(b4["total"]))),
         "эффективность по плановым / неплановым  задачам комендантов"),
        # Пара «плановая / неплановая», как в двух соседних плашках
        # (замечание 26.08.2026: неплановой уборки на слайде не хватало).
        (ru("%.1f%%/ %.1f%%" % (eff(b5["total"]), eff(b5b["total"]))),
         "эффективность по плановой / неплановой  уборке"),
    ]
    for i, (val, label) in enumerate(summary):
        x = 0.50 + (i % 3) * 4.02
        y = 2.00 + (i // 3) * 2.39
        D.round_rect(s, x, y, 3.78, 2.15, NAVY_DARK, radius=0.08)
        D.text(s, x + 0.25, y + 0.35, 3.28, 0.90, val,
               size=fit_size(val, 3.28, 30, 15, bold=True), bold=True, colorhex=BLUE, wrap=False)
        D.text(s, x + 0.25, y + 1.30, 3.28, 0.70, label, size=11.5, colorhex=WHITE, line_spacing=1.1)
    D.page += 1
    D.text(s, MARGIN_IN, 7.10, 8.00, 0.30, "Bruno · ВТБ · %s" % period_short,
           size=9, colorhex=TEXT_ON_NAVY_MUTED)
    D.text(s, SLIDE_W_IN - MARGIN_IN - 1, 7.10, 1.00, 0.30, str(D.page), size=9,
           colorhex=TEXT_ON_NAVY_MUTED, align=PP_ALIGN.RIGHT)

    # =================================================== 10. Разделитель приложений
    s = D.slide()
    D.rect(s, 0, 0, SLIDE_W_IN, SLIDE_H_IN, NAVY)
    D.text(s, MARGIN_IN, 2.90, 11.00, 0.50, "Приложения", size=16, bold=True, colorhex=BLUE)
    D.text(s, MARGIN_IN, 3.35, 11.30, 1.20, "Детализация и исходные данные", size=34, bold=True,
           colorhex=WHITE)
    D.text(s, MARGIN_IN, 4.35, 11.00, 1.00,
           "Помесячные таблицы и разбивка по объектам/комендантам для каждого блока отчёта",
           size=14, colorhex=TEXT_ON_NAVY_MUTED)
    D.footer(s, dark=True)

    # =================================================== приложения
    GRID_NOTE = "Мини-графики упрощены: подписаны только месяцы (без недельных отметок)."
    KOM_NOTE = ("Показаны коменданты с задачами ЭТОГО типа в текущем месяце и не менее %d "
                "задач за период; остальные учтены в общих цифрах слайда, но без личной строки."
                % KOM_MIN_TASKS)

    def grid_slide(num_label, title, subtitle, img_key, legend_items, note=None):
        pages = img.get(img_key) or []
        gx, gy, gw, gh = GRID_BOX
        for pi, (path, nrows) in enumerate(pages):
            # неполную страницу прижимаем к верху, иначе картинка центрируется
            # в боксе и над ней зияет пустая полоса
            box_h = min(gh, CELL_H * nrows)
            sl = D.slide()
            suffix = "" if len(pages) == 1 else "  ·  %d из %d" % (pi + 1, len(pages))
            D.text(sl, MARGIN_IN, 0.26, 11.30, 0.42, num_label + suffix, size=12, bold=True,
                   colorhex=BLUE)
            D.text(sl, MARGIN_IN, 0.54, 12.33, 0.42, title, size=19, bold=True, colorhex=NAVY,
                   wrap=False)
            D.text(sl, MARGIN_IN, 0.92, 12.33, 0.30, subtitle,
                   size=fit_size(subtitle, 12.33, 10.5, 8.2), colorhex=TEXT_MUTED, wrap=False)
            # Легенда — над сеткой: снизу место занято подписями месяцев и
            # помесячной эффективностью под каждым графиком.
            D.legend_row(sl, MARGIN_IN, 1.15, SLIDE_W_IN - 2 * MARGIN_IN,
                         legend_items, fontsize=8.5, row_h=0.2)
            D.picture(sl, path, gx, gy, gw, box_h)
            if note:
                D.text(sl, MARGIN_IN, 6.76, 12.33, 0.34, note, size=8.0, italic=True,
                       colorhex=TEXT_MUTED, line_spacing=1.1)
            D.footer(sl)

    def table_slide(num_label, title, subtitle, tkey, extra_rows=(), font_size=8.6,
                    col0_ratio=0.24, note=None, col_weights=None):
        t = tables.get(tkey)
        if not t or not t["rows"]:
            return
        sl = D.slide()
        D.text(sl, MARGIN_IN, 0.32, 11.30, 0.45, num_label, size=12, bold=True, colorhex=BLUE)
        D.text(sl, MARGIN_IN, 0.62, 12.33, 0.45, title, size=fit_size(title, 12.33, 19, 13),
               bold=True, colorhex=NAVY, wrap=False)
        D.text(sl, MARGIN_IN, 1.02, 12.33, 0.30, subtitle,
               size=fit_size(subtitle, 12.33, 10.5, 8.2), colorhex=TEXT_MUTED, wrap=False)
        rows = list(t["rows"]) + list(extra_rows)
        D.table(sl, MARGIN_IN, 1.42, SLIDE_W_IN - 2 * MARGIN_IN, 5.35 if not note else 5.10,
                t["header"], rows, col0_ratio=col0_ratio, font_size=font_size,
                col_weights=col_weights)
        if note:
            D.text(sl, MARGIN_IN, 6.84, 11.30, 0.30, note, size=8.3, italic=True, colorhex=TEXT_MUTED)
        D.footer(sl)

    def act_note(keep_set, what="объекты"):
        return ("Показаны %s с активностью ИМЕННО в этом блоке за текущий месяц — %d из %d "
                "в справочнике Bruno." % (what, len(keep_set), len(objs)))

    SUB_GRID = ("Выполнено / пропущено / в работе по неделям + эффективность (правая шкала "
                "своя у каждого графика), ряд стартует с первой активности объекта")
    SUB_GRID_KOM = ("Выполнено / пропущено / в работе по неделям + эффективность "
                    "(правая шкала своя у каждого графика)")
    SUB_TBL = ("Выполнено / пропущено / в работе по месяцам, эффективность каждого месяца "
               "и за весь период")
    SUB_TBL_KOM = "Выполнено / пропущено / в работе по месяцам, эффективность за весь период"
    # объём : эффективность = 3 : 1 на каждый месяц, последняя — итоговая эфф.
    SPLIT_W = [3.0, 1.05] * len(model["meta"]["months"]) + [1.25]

    table_slide("Приложение · Блок 2", "Обращения по объектам — помесячная динамика",
                "Количество обращений по QR, по месяцам", "b2",
                extra_rows=[tables["b2"]["total_row"]] if "b2" in tables else (), font_size=10)
    grid_slide("Приложение · Блок 3", "Все задачи — детализация по объектам",
               SUB_GRID, "b3_grid", GRID_LEGEND, act_note(keep_b3))
    table_slide("Приложение · Блок 3", "Все задачи — помесячная динамика по объектам",
                SUB_TBL, "b3", font_size=7.4, col0_ratio=0.17, note=act_note(keep_b3),
                col_weights=SPLIT_W)
    # Два новых слайда (замечание 26.08.2026): та же таблица, но отдельно по
    # плановым и по неплановым задачам — итоговая колонка «все задачи» скрывала,
    # за счёт чего именно получилась эффективность объекта.
    table_slide("Приложение · Блок 3", "Плановые задачи — помесячная динамика по объектам",
                SUB_TBL, "b3p", font_size=7.4, col0_ratio=0.17, note=act_note(keep_b3p),
                col_weights=SPLIT_W)
    table_slide("Приложение · Блок 3", "Неплановые задачи — помесячная динамика по объектам",
                SUB_TBL, "b3u", font_size=7.4, col0_ratio=0.17, note=act_note(keep_b3u),
                col_weights=SPLIT_W)
    grid_slide("Приложение · Блок 4", "Коменданты (неплановые) — детализация по каждому",
               SUB_GRID_KOM, "b4_grid", GRID_LEGEND, KOM_NOTE)
    table_slide("Приложение · Блок 4", "Коменданты (неплановые) — помесячная динамика",
                SUB_TBL_KOM, "b4", col0_ratio=0.20, note=KOM_NOTE)
    grid_slide("Приложение · Блок 4b", "Коменданты (плановые) — детализация по каждому",
               SUB_GRID_KOM, "b4b_grid", GRID_LEGEND, KOM_NOTE)
    table_slide("Приложение · Блок 4b", "Коменданты (плановые) — помесячная динамика",
                SUB_TBL_KOM, "b4b", col0_ratio=0.20, note=KOM_NOTE)
    grid_slide("Приложение · Блок 5", "Уборка (плановые) — детализация по объектам",
               SUB_GRID, "b5_grid", GRID_LEGEND, act_note(keep_b5))
    table_slide("Приложение · Блок 5", "Уборка (плановые) — помесячная динамика по объектам",
                SUB_TBL_KOM, "b5", font_size=8.2, col0_ratio=0.20, note=act_note(keep_b5))
    # Блок 5b добавлен по согласованию 22.08.2026 — детализация к новому графику
    # неплановой уборки на слайде 6. Состав ролей здесь шире, чем в блоке 5.
    clean_note = act_note(keep_b5b)
    if meta.get("clean_mode") == "split":
        clean_note = ("Учтены роли уборщиц и «Менеджер клининга» (в блоке 5 — только уборщицы); "
                      "объекты — с неплановой уборкой в текущем месяце, %d из %d."
                      % (len(keep_b5b), len(objs)))
    grid_slide("Приложение · Блок 5b", "Уборка (неплановые) — детализация по объектам",
               SUB_GRID, "b5b_grid", GRID_LEGEND, clean_note)
    table_slide("Приложение · Блок 5b", "Уборка (неплановые) — помесячная динамика по объектам",
                SUB_TBL_KOM, "b5b", font_size=8.2, col0_ratio=0.20, note=clean_note)

    D.prs.save(out_path)
    return len(D.prs.slides._sldIdLst)


def num_cell(v):
    """Число в таблицу по-русски: запятая вместо точки, прочерк вместо None."""
    return "—" if v is None else ru(str(v))

def short_num(n):
    n = float(n)
    if abs(n) >= 1_000_000:
        return ru("%.2f млн" % (n / 1_000_000))
    return fmt_n(n)


def tree_text(D, s, rect, model):
    """Нативный текст поверх tree_shapes.png. Геометрия продублирована из
    Charts.tree_shapes — при правке констант там править и здесь."""
    nx, ny, nw, nh = rect
    tree = model["tree"]
    X = lambda ax_x: nx + (ax_x / 100.0) * nw
    Y = lambda ax_y: ny + (1 - ax_y / 100.0) * nh
    WW = lambda ax_w: (ax_w / 100.0) * nw
    HH = lambda ax_h: (ax_h / 100.0) * nh

    root_x, root_y, root_w = 13, 50, 20
    D.text(s, X(root_x - root_w / 2), Y(root_y + 6), WW(root_w), HH(6), "Задач всего", size=12.5,
           bold=True, colorhex=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    D.text(s, X(root_x - root_w / 2), Y(root_y - 2), WW(root_w), HH(6), fmt_n(tree["total_due"]),
           size=11, colorhex=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    cats = [
        {"y": 74, "label": "Плановые", "data": tree["planned"], "extra": False, "mirror": False},
        # mirror сознательно выключен (п.19 списка правок): зеркальные подписи
        # осей у верхнего и нижнего блока читались как ошибка
        {"y": 26, "label": "Неплановые", "data": tree["unplanned"], "extra": True, "mirror": False},
    ]
    box_x, box_w, box_h = 60, 62, 36
    box_left, box_right = box_x - box_w / 2, box_x + box_w / 2
    spark_w, spark_h = box_w * 0.40, box_h * 0.52
    for c in cats:
        d = c["data"]
        tx = box_left + 3
        top = c["y"] + box_h / 2
        tw = box_w * 0.34
        D.text(s, X(tx), Y(top - 3.5), WW(tw), HH(4), c["label"], size=14.5, bold=True, colorhex=WHITE)
        D.text(s, X(tx), Y(top - 10), WW(tw), HH(3.5), "%s задач" % fmt_n(d["due"]), size=11, colorhex=WHITE)
        D.text(s, X(tx), Y(top - 15.5), WW(tw), HH(3.2), ru("эфф. %.1f%%" % d["efficiency"]),
               size=10, bold=True, colorhex=WHITE)
        if c["extra"]:
            note = ("Обращений: %s · неплановых задач по обращениям/заявкам: %s (%.0f%%)" % (
                fmt_n(tree["raw_feedback"]), fmt_n(tree["feedback_task_instances"]),
                tree["conversion_pct"]))
            D.text(s, X(tx), Y(top - 20.5), WW(tw), HH(13), note, size=8.0, colorhex=WHITE,
                   line_spacing=1.3)
        spark_left = box_right - spark_w - 3
        spark_top = c["y"] - spark_h / 2 + 1 + spark_h
        half = spark_w / 2
        count_label, eff_label = "Задач в месяц (столбцы)", "Эффективность, % (линия)"
        left_txt, right_txt = ((eff_label, count_label) if c["mirror"] else (count_label, eff_label))
        D.axis_caption(s, X(spark_left), Y(spark_top + 3.8), WW(half), left_txt, size=6.2,
                       colorhex=WHITE)
        D.axis_caption(s, X(spark_left + half), Y(spark_top + 3.8), WW(half), right_txt, size=6.2,
                       colorhex=WHITE, align=PP_ALIGN.RIGHT)


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--model", default="model.json.gz", help="model.json.gz конвейера 2.0")
    ap.add_argument("--out", default="Bruno_отчёт_ВТБ.pptx")
    ap.add_argument("--clean", choices=list(M.CLEAN_MODES), default="split",
                    help="кого считать уборщицами: base — только системная роль employee (как в "
                         "утверждённой колоде); base+mgr — вместе с «Менеджером клининга»; "
                         "split (по умолчанию, решение заказчика 22.08.2026) — плановые по "
                         "основной роли, неплановые по обеим")
    ap.add_argument("--months", type=int, default=3,
                    help="сколько ПОЛНЫХ месяцев брать до текущего (по умолчанию 3, т.е. горизонт "
                         "= текущий неполный + 3 = 4 месяца)")
    ap.add_argument("--from", dest="date_from", default=None, help="ГГГГ-ММ-ДД, перекрывает --months")
    ap.add_argument("--to", dest="date_to", default=None, help="ГГГГ-ММ-ДД")
    ap.add_argument("--imgdir", default=None, help="куда класть PNG (по умолчанию во временную папку)")
    args = ap.parse_args()

    raw = M.read_model(args.model)
    model = M.build(raw,
                    d_from=dt.date.fromisoformat(args.date_from) if args.date_from else None,
                    d_to=dt.date.fromisoformat(args.date_to) if args.date_to else None,
                    clean_mode=args.clean, months_back_n=args.months)

    imgdir = args.imgdir or tempfile.mkdtemp(prefix="bruno_v5_img_")
    print("=" * 66)
    print("Bruno -> презентация-копия утверждённой колоды")
    print("  Период:   %s .. %s (%d недель)" % (model["meta"]["period_from"],
                                                model["meta"]["period_to"], model["meta"]["n_weeks"]))
    print("  Уборщицы: %s (%s)" % (args.clean, model["meta"]["clean_note"]))
    print("  Объектов активных в текущем месяце:    %d" % len(model["active_objects_current_month"]))
    print("  Комендантов активных в текущем месяце: %d" % len(model["active_komendanty_current_month"]))
    print("=" * 66)

    n = build(model, args.out, imgdir)
    size_kb = os.path.getsize(args.out) / 1024
    print("  + готово: %s (%d слайдов, %.0f КБ)" % (args.out, n, size_kb))


if __name__ == "__main__":
    main()
