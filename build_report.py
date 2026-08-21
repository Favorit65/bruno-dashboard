#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Сборка отчёта-презентации (PPTX) по данным дашборда 2.0.

Вход  — model.json.gz (собирается build_model.py из архива Bruno).
Выход — .pptx на 8 слайдов за произвольный период.

    python build_report.py --model model.json.gz --from 2026-07-01 --to 2026-07-31

Оформление повторяет утверждённый отчёт (БРУНО_исправленно.pptx, июль 2026):
палитра ВТБ, тёмно-синие титул и «Итоги», светлые контентные слайды с
плитками KPI и графиками.

Что НЕ вошло в эту версию (осознанно, см. ПРОГРЕСС_автоматизация_git_и_отчёт.md):
слайд оклейки QR (данные ведутся вручную, их нет в API), приложения с
детализацией по каждому коменданту и скорость уборки по башням (нужны поля,
которых нет в агрегате model.json.gz — только в сыром архиве).
"""

import argparse
import gzip
import json
import os
import sys
import tempfile
from collections import defaultdict
from datetime import date, timedelta

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.ticker import FuncFormatter

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# --------------------------------------------------------------- палитра ВТБ
NAVY = "002882"
NAVY_DARK = "001A57"
BLUE = "00AAFF"
GREEN = "23A75B"
RED = "E62632"
AMBER = "F5A623"
WHITE = "FFFFFF"
BG_TILE = "F4F6FA"
GRID = "E2E7F0"
TEXT_DARK = "1A2233"
TEXT_MUTED = "5B6B82"
TEXT_ON_NAVY_MUTED = "A9BEDD"

FONT = "Golos Text"          # корпоративный шрифт ВТБ
FONT_MPL = "DejaVu Sans"     # для картинок графиков

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.5

# единая сетка контентного слайда: шапка -> плитки KPI -> график -> легенда -> сноска
CHART_X = 1.1
CHART_Y = 3.15
CHART_W = 11.1
CHART_H = 3.0
LEGEND_Y = 6.28
NOTE_Y = 6.62

STATUSES = ("NEW", "WAITING", "COMPLETING", "COMPLETED", "MISSED")
MONTHS_RU = ["", "январь", "февраль", "март", "апрель", "май", "июнь", "июль",
             "август", "сентябрь", "октябрь", "ноябрь", "декабрь"]
MONTHS_RU_SHORT = ["", "янв", "фев", "мар", "апр", "май", "июн", "июл",
                   "авг", "сен", "окт", "ноя", "дек"]


# --------------------------------------------------------------- утилиты
def rgb(hex6):
    return RGBColor.from_string(hex6)


def fmt_int(n):
    return f"{int(round(n)):,}".replace(",", " ")


def fmt_short(n):
    """1 058 636 -> «1,06 млн», 43 833 -> «43 833» (как в утверждённом отчёте)."""
    n = float(n)
    if abs(n) >= 1_000_000:
        return f"{n / 1_000_000:.2f}".replace(".", ",") + " млн"
    if abs(n) >= 100_000:
        return f"{n / 1000:.0f}".replace(".", ",") + " тыс."
    return fmt_int(n)


def fmt_pct(x, digits=1):
    if x is None:
        return "—"
    return f"{x * 100:.{digits}f}".replace(".", ",") + "%"


def eff(c, m):
    den = c + m
    return c / den if den else None


def new_bucket():
    return {k: 0 for k in STATUSES} | {"completedOnTime": 0, "completedLate": 0,
                                       "delayHoursSum": 0.0, "delayHoursN": 0}


def acc(dst, src):
    for k in STATUSES + ("completedOnTime", "completedLate", "delayHoursN"):
        dst[k] += src.get(k, 0)
    dst["delayHoursSum"] += src.get("delayHoursSum", 0.0)


def btotal(b):
    return sum(b[k] for k in STATUSES)


def daterange(d0, d1):
    d = d0
    while d <= d1:
        yield d
        d += timedelta(days=1)


def month_label(y, m, short=False):
    return (MONTHS_RU_SHORT[m] if short else MONTHS_RU[m]).capitalize() + (f" {y}" if not short else "")


# --------------------------------------------------------------- загрузка
def _read_json(path):
    op = gzip.open if str(path).endswith(".gz") else open
    with op(path, "rt", encoding="utf-8") as f:
        return json.load(f)


def load_model(path):
    """Читает модель формата 1 (model.json.gz) или формата 2 (model.core.json.gz).

    Формат 2 разрезан на две части ради скорости загрузки сайта; отчёту нужны обе,
    поэтому вторая (model.cubes.json.gz) подхватывается рядом с первой. На выходе
    в обоих случаях — одна и та же структура, дальше по коду разницы нет.
    """
    data = _read_json(path)
    if not isinstance(data, dict) or data.get("f") != 2:
        return data

    import model_v2
    from pathlib import Path as _P
    p = _P(path)
    name = p.name
    cubes_name = (name[:-len(".core.json.gz")] + ".cubes.json.gz"
                  if name.endswith(".core.json.gz") else "model.cubes.json.gz")
    cubes_path = p.with_name(cubes_name)
    if not cubes_path.exists():
        raise SystemExit("ОШИБКА: рядом с %s нет второй половины модели (%s)" % (name, cubes_name))
    return model_v2.decode(data, _read_json(cubes_path))


class Period:
    """Агрегаты за выбранный период, посчитанные из кубов model.json.gz."""

    def __init__(self, model, d_from, d_to):
        self.model = model
        self.d_from = d_from
        self.d_to = d_to
        self.objects = model["directories"]["objects"]
        self.days = [d for d in daterange(d_from, d_to)]
        self.in_range = {d.isoformat() for d in self.days}

        self.planned = new_bucket()
        self.unplanned = new_bucket()
        self.by_object = defaultdict(lambda: {"planned": new_bucket(), "unplanned": new_bucket()})
        self.by_day = defaultdict(lambda: {"planned": new_bucket(), "unplanned": new_bucket()})
        self.by_role = defaultdict(lambda: {"planned": new_bucket(), "unplanned": new_bucket()})
        self.by_hour = defaultdict(lambda: {"planned": new_bucket(), "unplanned": new_bucket()})

        for key, val in model["daily"]["byObjectDay"].items():
            ds, oid = key.split("|", 1)
            if ds not in self.in_range:
                continue
            for kind, tgt in (("planned", self.planned), ("unplanned", self.unplanned)):
                acc(tgt, val[kind])
                acc(self.by_object[oid][kind], val[kind])
                acc(self.by_day[ds][kind], val[kind])

        for key, val in model["daily"].get("byObjectRoleDay", {}).items():
            parts = key.split("|")
            if len(parts) != 3 or parts[0] not in self.in_range:
                continue
            role = parts[2].strip() or "Без роли"
            for kind in ("planned", "unplanned"):
                acc(self.by_role[role][kind], val[kind])

        for key, val in model["daily"].get("byObjectHour", {}).items():
            parts = key.split("|")
            if len(parts) != 3 or parts[0] not in self.in_range:
                continue
            hour = int(parts[1])
            for kind in ("planned", "unplanned"):
                acc(self.by_hour[hour][kind], val[kind])

    # --- сводные показатели -------------------------------------------------
    @property
    def total(self):
        t = new_bucket()
        acc(t, self.planned)
        acc(t, self.unplanned)
        return t

    def eff_of(self, b):
        return eff(b["COMPLETED"], b["MISSED"])

    @property
    def active_objects(self):
        return [oid for oid, v in self.by_object.items()
                if btotal(v["planned"]) + btotal(v["unplanned"]) > 0]

    def buckets(self):
        """Разбивка периода на недели или месяцы — в зависимости от длины."""
        n = len(self.days)
        if n <= 62:
            out, cur = [], None
            for d in self.days:
                ws = d - timedelta(days=d.weekday())
                if cur is None or cur["key"] != ws:
                    # подписываем неделю ПЕРВЫМ днём, попавшим в период, а не
                    # понедельником календарной недели — иначе первая колонка
                    # подписана датой, которой в отчёте нет
                    cur = {"key": ws, "label": d.strftime("%d.%m"), "days": []}
                    out.append(cur)
                cur["days"].append(d)
            return out, "неделям"
        out, cur = [], None
        for d in self.days:
            mk = (d.year, d.month)
            if cur is None or cur["key"] != mk:
                cur = {"key": mk, "label": MONTHS_RU_SHORT[d.month].capitalize(), "days": []}
                out.append(cur)
            cur["days"].append(d)
        return out, "месяцам"

    def bucket_series(self):
        buckets, gran = self.buckets()
        rows = []
        for b in buckets:
            agg = {"planned": new_bucket(), "unplanned": new_bucket()}
            for d in b["days"]:
                v = self.by_day.get(d.isoformat())
                if not v:
                    continue
                acc(agg["planned"], v["planned"])
                acc(agg["unplanned"], v["unplanned"])
            tot = new_bucket()
            acc(tot, agg["planned"])
            acc(tot, agg["unplanned"])
            rows.append({"label": b["label"], "planned": agg["planned"],
                         "unplanned": agg["unplanned"], "total": tot})
        return rows, gran

    def objects_sorted(self):
        rows = []
        for oid in self.active_objects:
            v = self.by_object[oid]
            tot = new_bucket()
            acc(tot, v["planned"])
            acc(tot, v["unplanned"])
            rows.append({
                "id": oid,
                "name": self.objects.get(oid) or ("Объект " + oid[:8]),
                "total": btotal(tot),
                "completed": tot["COMPLETED"],
                "missed": tot["MISSED"],
                "eff": self.eff_of(tot),
                "unplanned": btotal(v["unplanned"]),
            })
        rows.sort(key=lambda r: -r["total"])
        return rows

    def roles_sorted(self):
        rows = []
        for role, v in self.by_role.items():
            tot = new_bucket()
            acc(tot, v["planned"])
            acc(tot, v["unplanned"])
            if btotal(tot) <= 0:
                continue
            rows.append({
                "role": role,
                "total": btotal(tot),
                "planned": btotal(v["planned"]),
                "unplanned": btotal(v["unplanned"]),
                "eff": self.eff_of(tot),
            })
        rows.sort(key=lambda r: -r["total"])
        return rows


# --------------------------------------------------------------- графики
def mpl_setup():
    plt.rcParams.update({
        "font.family": FONT_MPL,
        "axes.edgecolor": "#" + GRID,
        "axes.labelcolor": "#" + TEXT_MUTED,
        "xtick.color": "#" + TEXT_MUTED,
        "ytick.color": "#" + TEXT_MUTED,
        "text.color": "#" + TEXT_DARK,
        "figure.dpi": 200,
        "savefig.transparent": True,
    })


def _thousands(x, _):
    return fmt_int(x) if x else "0"


def chart_dynamics(rows, path, w=11.0, h=3.0):
    """Столбцы выполнено/пропущено + линия эффективности на правой оси."""
    fig, ax = plt.subplots(figsize=(w, h))
    xs = range(len(rows))
    done = [r["total"]["COMPLETED"] for r in rows]
    miss = [r["total"]["MISSED"] for r in rows]
    work = [r["total"]["COMPLETING"] + r["total"]["NEW"] + r["total"]["WAITING"] for r in rows]
    ax.bar(xs, done, color="#" + GREEN, width=0.62, label="Выполнено")
    ax.bar(xs, miss, bottom=done, color="#" + RED, width=0.62, label="Пропущено")
    ax.bar(xs, work, bottom=[d + m for d, m in zip(done, miss)], color="#" + AMBER,
           width=0.62, label="В работе")
    ax.set_xticks(list(xs))
    ax.set_xticklabels([r["label"] for r in rows], fontsize=7.5)
    ax.yaxis.set_major_formatter(FuncFormatter(_thousands))
    ax.tick_params(axis="y", labelsize=7.5)
    ax.grid(axis="y", color="#" + GRID, linewidth=0.6)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)

    ax2 = ax.twinx()
    effs = [(eff(r["total"]["COMPLETED"], r["total"]["MISSED"]) or 0) * 100 for r in rows]
    ax2.plot(list(xs), effs, color="#" + NAVY, linewidth=2.0, marker="o", markersize=4)
    for x, e in zip(xs, effs):
        ax2.annotate(f"{e:.0f}%", (x, e), textcoords="offset points", xytext=(0, 7),
                     ha="center", fontsize=7.5, color="#" + NAVY, fontweight="bold")
    ax2.set_ylim(0, max(effs + [1]) * 1.45)
    ax2.set_yticks([])
    for s in ("top", "right", "left", "bottom"):
        ax2.spines[s].set_visible(False)

    fig.tight_layout(pad=0.2)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def chart_objects(rows, path, w=11.0, h=3.15, top=10):
    """Горизонтальные полосы: объём задач по объектам + подпись эффективности."""
    rows = rows[:top][::-1]
    fig, ax = plt.subplots(figsize=(w, h))
    ys = range(len(rows))
    done = [r["completed"] for r in rows]
    miss = [r["missed"] for r in rows]
    other = [r["total"] - r["completed"] - r["missed"] for r in rows]
    ax.barh(ys, done, color="#" + GREEN, height=0.62, label="Выполнено")
    ax.barh(ys, miss, left=done, color="#" + RED, height=0.62, label="Пропущено")
    ax.barh(ys, other, left=[d + m for d, m in zip(done, miss)], color="#" + AMBER,
            height=0.62, label="В работе")
    ax.set_yticks(list(ys))
    ax.set_yticklabels([r["name"][:34] for r in rows], fontsize=8)
    ax.xaxis.set_major_formatter(FuncFormatter(_thousands))
    ax.tick_params(axis="x", labelsize=7.5)
    ax.grid(axis="x", color="#" + GRID, linewidth=0.6)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    xmax = max([r["total"] for r in rows] + [1])
    ax.set_xlim(0, xmax * 1.16)
    for y, r in zip(ys, rows):
        ax.annotate(f"{fmt_short(r['total'])} · {fmt_pct(r['eff'], 0)}",
                    (r["total"], y), xytext=(6, 0), textcoords="offset points",
                    va="center", fontsize=7.5, color="#" + TEXT_MUTED)
    fig.tight_layout(pad=0.2)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def chart_roles(rows, path, w=11.0, h=3.15, top=8):
    rows = rows[:top][::-1]
    fig, ax = plt.subplots(figsize=(w, h))
    ys = range(len(rows))
    pl = [r["planned"] for r in rows]
    un = [r["unplanned"] for r in rows]
    ax.barh(ys, pl, color="#" + NAVY, height=0.58, label="Плановые")
    ax.barh(ys, un, left=pl, color="#" + BLUE, height=0.58, label="Внеплановые")
    ax.set_yticks(list(ys))
    ax.set_yticklabels([r["role"][:28] for r in rows], fontsize=8.5)
    ax.xaxis.set_major_formatter(FuncFormatter(_thousands))
    ax.tick_params(axis="x", labelsize=7.5)
    ax.grid(axis="x", color="#" + GRID, linewidth=0.6)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    xmax = max([r["total"] for r in rows] + [1])
    ax.set_xlim(0, xmax * 1.18)
    for y, r in zip(ys, rows):
        ax.annotate(f"{fmt_short(r['total'])} · {fmt_pct(r['eff'], 0)}",
                    (r["total"], y), xytext=(6, 0), textcoords="offset points",
                    va="center", fontsize=7.5, color="#" + TEXT_MUTED)
    fig.tight_layout(pad=0.2)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


def chart_hours(by_hour, path, w=11.0, h=2.95):
    fig, ax = plt.subplots(figsize=(w, h))
    hours = list(range(24))
    pl, un = [], []
    for hh in hours:
        v = by_hour.get(hh)
        pl.append(btotal(v["planned"]) if v else 0)
        un.append(btotal(v["unplanned"]) if v else 0)
    ax.bar(hours, pl, color="#" + NAVY, width=0.7, label="Плановые")
    ax.bar(hours, un, bottom=pl, color="#" + BLUE, width=0.7, label="Внеплановые")
    ax.set_xticks(hours)
    ax.set_xticklabels([f"{h:02d}" for h in hours], fontsize=7)
    ax.yaxis.set_major_formatter(FuncFormatter(_thousands))
    ax.tick_params(axis="y", labelsize=7.5)
    ax.grid(axis="y", color="#" + GRID, linewidth=0.6)
    ax.set_axisbelow(True)
    for s in ("top", "right", "left"):
        ax.spines[s].set_visible(False)
    fig.tight_layout(pad=0.2)
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)


# --------------------------------------------------------------- вёрстка
class Deck:
    def __init__(self, footer):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self.footer = footer
        self.page = 0

    # --- примитивы ---
    def _blank(self):
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def rect(self, slide, x, y, w, h, fill, shape=MSO_SHAPE.RECTANGLE):
        sh = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = rgb(fill)
        sh.line.fill.background()
        sh.shadow.inherit = False
        return sh

    def text(self, slide, x, y, w, h, s, size=12, color=TEXT_DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        lines = str(s).split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            p.line_spacing = spacing
            r = p.add_run()
            r.text = line
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.name = FONT
            r.font.color.rgb = rgb(color)
        return tb

    def footer_line(self, slide, dark=False):
        self.page += 1
        col = TEXT_ON_NAVY_MUTED if dark else TEXT_MUTED
        self.text(slide, MARGIN, SLIDE_H - 0.38, 8.0, 0.3, self.footer, size=9, color=col)
        self.text(slide, SLIDE_W - 1.5, SLIDE_H - 0.38, 1.0, 0.3, str(self.page),
                  size=9, color=col, align=PP_ALIGN.RIGHT)

    # --- типы слайдов ---
    def title_slide(self, title, subtitle, period_str, note):
        s = self._blank()
        self.rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
        self.rect(s, 0, 0, 0.22, SLIDE_H, BLUE)
        self.text(s, 1.1, 2.25, 10, 0.5, "BRUNO", size=15, color=BLUE, bold=True)
        self.text(s, 1.1, 2.75, 11, 1.0, title, size=40, color=WHITE, bold=True)
        self.text(s, 1.1, 3.95, 11, 0.5, subtitle, size=14, color=TEXT_ON_NAVY_MUTED)
        self.rect(s, 1.1, 4.65, 2.6, 0.035, BLUE)
        self.text(s, 1.1, 4.95, 11, 0.4, period_str, size=17, color=WHITE, bold=True)
        self.text(s, 1.1, 5.45, 11, 0.4, note, size=11, color=TEXT_ON_NAVY_MUTED)
        self.page += 1
        return s

    def content_slide(self, num, title, subtitle):
        s = self._blank()
        badge = self.rect(s, MARGIN, 0.42, 0.5, 0.5, NAVY)
        tf = badge.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(num)
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.name = FONT
        r.font.color.rgb = rgb(WHITE)
        self.text(s, 1.18, 0.35, 10.5, 0.5, title, size=22, color=NAVY, bold=True)
        self.text(s, 1.18, 0.86, 11.0, 0.4, subtitle, size=12.5, color=TEXT_MUTED)
        return s

    def tiles(self, slide, items, y=1.62, x0=MARGIN, w=2.9, h=1.15, gap=0.24):
        """Светлые плитки KPI: [(значение, подпись, цвет_значения), ...]"""
        for i, (val, lbl, col) in enumerate(items):
            x = x0 + i * (w + gap)
            self.rect(slide, x, y, w, h, BG_TILE)
            # длинные значения (например название объекта) уменьшаем, чтобы не
            # переносились на вторую строку и не вылезали из плитки
            n = len(str(val))
            size = 24 if n <= 11 else (19 if n <= 15 else (16 if n <= 20 else 13))
            self.text(slide, x + 0.16, y + 0.16, w - 0.3, 0.55, val, size=size, color=col, bold=True)
            self.text(slide, x + 0.16, y + 0.74, w - 0.3, 0.36, lbl, size=10.5, color=TEXT_MUTED)

    def fit_picture(self, slide, png, x, y, w, h):
        """Вписывает картинку в прямоугольник с сохранением пропорций, по центру."""
        from PIL import Image
        iw, ih = Image.open(png).size
        scale = min(w / iw, h / ih)
        pw, ph = iw * scale, ih * scale
        slide.shapes.add_picture(png, Inches(x + (w - pw) / 2), Inches(y + (h - ph) / 2),
                                 width=Inches(pw), height=Inches(ph))

    def legend(self, slide, items, x, y):
        cur = x
        for label, color in items:
            self.rect(slide, cur, y, 0.13, 0.13, color)
            self.text(slide, cur + 0.21, y - 0.01, 1.3, 0.18, label, size=8.7, color=TEXT_DARK)
            cur += 0.21 + 0.11 * len(label) + 0.18

    def summary_slide(self, title, subtitle, items):
        """Тёмный слайд «Итоги»: 6 плиток 3x2."""
        s = self._blank()
        self.rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
        self.text(s, MARGIN, 0.5, 11, 0.6, title, size=28, color=WHITE, bold=True)
        self.text(s, MARGIN, 1.15, 11, 0.4, subtitle, size=14, color=TEXT_ON_NAVY_MUTED)
        w, h, gap = 3.78, 2.15, 0.24
        for i, (val, lbl) in enumerate(items[:6]):
            col_i, row_i = i % 3, i // 3
            x = MARGIN + col_i * (w + gap)
            y = 2.0 + row_i * (h + gap)
            self.rect(s, x, y, w, h, NAVY_DARK)
            self.text(s, x + 0.25, y + 0.35, w - 0.5, 0.9, val, size=30, color=BLUE, bold=True)
            self.text(s, x + 0.25, y + 1.3, w - 0.5, 0.7, lbl, size=11.5, color=WHITE, spacing=1.15)
        self.footer_line(s, dark=True)
        return s

    def table(self, slide, x, y, w, h, headers, rows, col_widths=None,
              align_right_from=1, font=9.0):
        shape = slide.shapes.add_table(len(rows) + 1, len(headers),
                                       Inches(x), Inches(y), Inches(w), Inches(h))
        tbl = shape.table
        tbl.first_row = True
        if col_widths:
            total = sum(col_widths)
            for i, cw in enumerate(col_widths):
                tbl.columns[i].width = Inches(w * cw / total)
        for j, htxt in enumerate(headers):
            c = tbl.cell(0, j)
            c.text = ""
            c.fill.solid()
            c.fill.fore_color.rgb = rgb(NAVY)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT if j >= align_right_from else PP_ALIGN.LEFT
            r = p.add_run()
            r.text = htxt
            r.font.size = Pt(font)
            r.font.bold = True
            r.font.name = FONT
            r.font.color.rgb = rgb(WHITE)
        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                c = tbl.cell(i, j)
                c.text = ""
                c.fill.solid()
                c.fill.fore_color.rgb = rgb(WHITE if i % 2 else BG_TILE)
                c.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = c.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.RIGHT if j >= align_right_from else PP_ALIGN.LEFT
                r = p.add_run()
                r.text = str(val)
                r.font.size = Pt(font)
                r.font.name = FONT
                r.font.color.rgb = rgb(TEXT_DARK)
        for i in range(len(rows) + 1):
            tbl.rows[i].height = Inches(h / (len(rows) + 1))
        return tbl


# --------------------------------------------------------------- сборка
def build(model, d_from, d_to, out_path, title="Отчёт по эксплуатации объектов"):
    mpl_setup()
    P = Period(model, d_from, d_to)
    period_str = f"{d_from.strftime('%d.%m.%Y')} — {d_to.strftime('%d.%m.%Y')}"
    period_short = f"{d_from.strftime('%d.%m')}–{d_to.strftime('%d.%m.%Y')}"
    footer = f"Bruno · ВТБ · {period_short}"
    deck = Deck(footer)

    tot = P.total
    n_total = btotal(tot)
    n_planned = btotal(P.planned)
    n_unplanned = btotal(P.unplanned)
    e_all = P.eff_of(tot)
    e_pl = P.eff_of(P.planned)
    e_un = P.eff_of(P.unplanned)
    objects = P.objects_sorted()
    roles = P.roles_sorted()
    series, gran = P.bucket_series()

    snap = model.get("meta", {}).get("generatedFromSnapshot", "—")
    tmpdir = tempfile.mkdtemp(prefix="bruno_report_")

    # --- 1. титул ---
    deck.title_slide(
        title,
        "Плановые и внеплановые задачи · объекты · роли · нагрузка",
        period_str,
        f"Источник: система Bruno, архив на {snap}. Данные и методика — те же, "
        f"что в дашборде.",
    )

    # --- 2. портфель задач ---
    s = deck.content_slide(1, "Портфель задач за период",
                           "Сколько задач заведено, как они распределены и что из них выполнено")
    deck.tiles(s, [
        (fmt_short(n_total), "задач всего", NAVY),
        (fmt_pct(e_all), "эффективность за период", RED if (e_all or 0) < 0.5 else GREEN),
        (fmt_short(tot["COMPLETED"]), "выполнено", NAVY),
        (fmt_short(n_unplanned), "внеплановых задач", NAVY),
    ])
    png = os.path.join(tmpdir, "dyn.png")
    chart_dynamics(series, png)
    deck.text(s, CHART_X, 2.93, 4.47, 0.2, f"Задач в {'неделю' if gran == 'неделям' else 'месяц'}",
              size=8.7, color=TEXT_MUTED)
    deck.text(s, CHART_X + CHART_W - 4.47, 2.93, 4.47, 0.2, "Эффективность, %", size=8.7,
              color=NAVY, align=PP_ALIGN.RIGHT)
    deck.fit_picture(s, png, CHART_X, CHART_Y, CHART_W, CHART_H)
    deck.legend(s, [("Выполнено", GREEN), ("Пропущено", RED), ("В работе", AMBER)],
                4.0, LEGEND_Y)
    deck.text(s, MARGIN, NOTE_Y, 12.3, 0.4,
              f"Плановые — {fmt_short(n_planned)} задач ({fmt_pct(e_pl)} выполнено), "
              f"внеплановые (обращения и заявки) — {fmt_short(n_unplanned)} "
              f"({fmt_pct(e_un)}). Эффективность = выполнено / (выполнено + пропущено).",
              size=8.6, color=TEXT_MUTED)
    deck.footer_line(s)

    # --- 3. объекты ---
    s = deck.content_slide(2, "Объекты", "Объём задач и эффективность по объектам портфеля")
    top_obj = objects[0] if objects else None
    deck.tiles(s, [
        (str(len(objects)), "активных объектов", NAVY),
        (top_obj["name"][:18] if top_obj else "—", "наибольший объём", NAVY),
        (fmt_pct(max((o["eff"] or 0) for o in objects)) if objects else "—",
         "лучшая эффективность", GREEN),
        (fmt_pct(min((o["eff"] or 0) for o in objects)) if objects else "—",
         "худшая эффективность", RED),
    ])
    png = os.path.join(tmpdir, "obj.png")
    chart_objects(objects, png)
    deck.fit_picture(s, png, CHART_X, CHART_Y - 0.15, CHART_W, CHART_H + 0.15)
    deck.legend(s, [("Выполнено", GREEN), ("Пропущено", RED), ("В работе", AMBER)],
                4.0, LEGEND_Y)
    deck.text(s, MARGIN, NOTE_Y, 12.3, 0.3,
              "Показаны 10 объектов с наибольшим объёмом задач. Справа от полосы — "
              "всего задач и эффективность.", size=8.6, color=TEXT_MUTED)
    deck.footer_line(s)

    # --- 4. роли ---
    s = deck.content_slide(3, "Персонал по ролям",
                           "Сколько задач приходится на каждую роль и как они закрываются")
    if roles:
        top_role = roles[0]
        deck.tiles(s, [
            (str(len(roles)), "ролей в работе", NAVY),
            (top_role["role"][:18], "основная нагрузка", NAVY),
            (fmt_short(top_role["total"]), "задач у неё", NAVY),
            (fmt_pct(top_role["eff"]), "её эффективность",
             RED if (top_role["eff"] or 0) < 0.5 else GREEN),
        ])
    png = os.path.join(tmpdir, "roles.png")
    chart_roles(roles, png)
    deck.fit_picture(s, png, CHART_X, CHART_Y - 0.15, CHART_W, CHART_H + 0.15)
    deck.legend(s, [("Плановые", NAVY), ("Внеплановые", BLUE)], 4.8, LEGEND_Y)
    deck.text(s, MARGIN, NOTE_Y, 12.3, 0.3,
              "Роль задачи определяется по команде исполнителя; для внеплановых задач без "
              "команды — по назначенным сотрудникам.", size=8.6, color=TEXT_MUTED)
    deck.footer_line(s)

    # --- 5. нагрузка по часам ---
    s = deck.content_slide(4, "Нагрузка по часам суток",
                           "Когда задачи должны быть выполнены — распределение по времени")
    hours_tot = {h: btotal(v["planned"]) + btotal(v["unplanned"]) for h, v in P.by_hour.items()}
    peak = max(hours_tot, key=hours_tot.get) if hours_tot else None
    night = sum(v for h, v in hours_tot.items() if h < 6 or h >= 22)
    deck.tiles(s, [
        (f"{peak:02d}:00" if peak is not None else "—", "пиковый час", NAVY),
        (fmt_short(hours_tot.get(peak, 0)) if peak is not None else "—", "задач в пиковый час", NAVY),
        (fmt_pct(night / sum(hours_tot.values()) if hours_tot else None, 1), "в ночные часы", NAVY),
        (fmt_short(sum(hours_tot.values()) / max(len(P.days), 1)), "задач в среднем за день", NAVY),
    ])
    png = os.path.join(tmpdir, "hours.png")
    chart_hours(P.by_hour, png)
    deck.fit_picture(s, png, CHART_X, CHART_Y, CHART_W, CHART_H)
    deck.legend(s, [("Плановые", NAVY), ("Внеплановые", BLUE)], 4.8, LEGEND_Y)
    deck.text(s, MARGIN, NOTE_Y, 12.3, 0.4,
              "Час определяется по плановому времени выполнения задачи (МСК). "
              "Ночные часы — с 22:00 до 06:00.", size=8.6, color=TEXT_MUTED)
    deck.footer_line(s)

    # --- 6. итоги (тёмный) ---
    deck.summary_slide("Итоги периода", period_str, [
        (fmt_short(n_total), "задач заведено в Bruno"),
        (fmt_pct(e_all), "выполнено от суммы «выполнено + пропущено»"),
        (str(len(objects)), "объектов с активностью за период"),
        (f"{fmt_pct(e_pl)} / {fmt_pct(e_un)}", "эффективность по плановым / внеплановым задачам"),
        (fmt_short(n_unplanned), "внеплановых задач (обращения и заявки)"),
        (fmt_short(n_total / max(len(P.days), 1)), "задач в среднем за сутки"),
    ])

    # --- 7. приложение: таблица по объектам ---
    s = deck.content_slide(5, "Приложение · объекты",
                           "Полная детализация по всем объектам портфеля за период")
    rows = [[o["name"][:40], fmt_int(o["total"]), fmt_int(o["completed"]),
             fmt_int(o["missed"]), fmt_int(o["unplanned"]), fmt_pct(o["eff"])]
            for o in objects]
    if rows:
        h = min(5.4, 0.32 * (len(rows) + 1))
        deck.table(s, MARGIN, 1.55, SLIDE_W - 2 * MARGIN, h,
                   ["Объект", "Всего", "Выполнено", "Пропущено", "Внеплановых", "Эффективность"],
                   rows, col_widths=[4.2, 1.2, 1.4, 1.3, 1.4, 1.5],
                   font=9.0 if len(rows) <= 14 else 8.0)
    deck.text(s, MARGIN, SLIDE_H - 0.72, 12.3, 0.3,
              "Эффективность = выполнено / (выполнено + пропущено). Задачи со статусами "
              "«новая», «ожидает», «в работе» в знаменатель не входят.",
              size=8.6, color=TEXT_MUTED)
    deck.footer_line(s)

    deck.prs.save(out_path)
    return {
        "slides": len(deck.prs.slides.__iter__.__self__._sldIdLst),
        "total": n_total,
        "objects": len(objects),
        "roles": len(roles),
        "eff": e_all,
    }


def main():
    ap = argparse.ArgumentParser(description=__doc__,
                                 formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("--model", default="model.json.gz", help="Путь к model.json.gz")
    ap.add_argument("--from", dest="date_from", help="Начало периода, ГГГГ-ММ-ДД")
    ap.add_argument("--to", dest="date_to", help="Конец периода, ГГГГ-ММ-ДД")
    ap.add_argument("--preset", choices=["last-month", "last-30", "last-90", "ytd", "all"],
                    help="Готовый период вместо --from/--to")
    ap.add_argument("--out", default="report.pptx", help="Имя выходного файла")
    ap.add_argument("--title", default="Отчёт по эксплуатации объектов")
    args = ap.parse_args()

    model = load_model(args.model)
    meta = model.get("meta", {})
    a_from = date.fromisoformat(meta.get("archiveFrom", "2026-01-01"))
    a_to = date.fromisoformat(meta.get("archiveTo", date.today().isoformat()))

    if args.preset:
        if args.preset == "all":
            d_from, d_to = a_from, a_to
        elif args.preset == "ytd":
            d_from, d_to = date(a_to.year, 1, 1), a_to
        elif args.preset == "last-month":
            first_this = a_to.replace(day=1)
            d_to = first_this - timedelta(days=1)
            d_from = d_to.replace(day=1)
        elif args.preset == "last-30":
            d_from, d_to = a_to - timedelta(days=29), a_to
        else:
            d_from, d_to = a_to - timedelta(days=89), a_to
    else:
        d_from = date.fromisoformat(args.date_from) if args.date_from else a_from
        d_to = date.fromisoformat(args.date_to) if args.date_to else a_to

    d_from = max(d_from, a_from)
    d_to = min(d_to, a_to)
    if d_from > d_to:
        sys.exit(f"ОШИБКА: пустой период {d_from}..{d_to} (архив: {a_from}..{a_to})")

    print("=" * 66)
    print("Bruno -> отчёт-презентация")
    print(f"  Период:  {d_from} .. {d_to}  ({(d_to - d_from).days + 1} дн.)")
    print(f"  Архив:   {a_from} .. {a_to}")
    print(f"  Выход:   {args.out}")
    print("=" * 66)

    info = build(model, d_from, d_to, args.out, title=args.title)
    size_kb = os.path.getsize(args.out) / 1024
    print(f"  + готово: {args.out} ({size_kb:,.0f} КБ)")
    print(f"    задач за период: {fmt_int(info['total'])}, объектов: {info['objects']}, "
          f"ролей: {info['roles']}, эффективность: {fmt_pct(info['eff'])}")


if __name__ == "__main__":
    main()
